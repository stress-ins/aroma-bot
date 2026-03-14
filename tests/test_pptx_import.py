"""Tests for PPTX import — extract_images_from_pptx and carousel import endpoint."""
from __future__ import annotations

import io
from unittest.mock import AsyncMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

from bot.services.carousel_assets import extract_images_from_pptx


def _make_test_pptx(slide_count: int = 3, with_images: bool = True) -> bytes:
    """Create a minimal PPTX with optional images for testing."""
    prs = Presentation()
    # Create a small PNG (1x1 pixel, red)
    png_bytes = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    for _ in range(slide_count):
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        if with_images:
            slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(1), Inches(1), Inches(4), Inches(3))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# extract_images_from_pptx
# ---------------------------------------------------------------------------

def test_extract_images_basic():
    pptx_bytes = _make_test_pptx(slide_count=3, with_images=True)
    images = extract_images_from_pptx(pptx_bytes)
    assert len(images) == 3
    assert all(img is not None for img in images)
    # Each extracted image should be valid PNG bytes
    for img in images:
        assert img[:4] == b"\x89PNG"


def test_extract_images_no_images():
    pptx_bytes = _make_test_pptx(slide_count=2, with_images=False)
    images = extract_images_from_pptx(pptx_bytes)
    assert len(images) == 2
    assert all(img is None for img in images)


def test_extract_images_single_slide():
    pptx_bytes = _make_test_pptx(slide_count=1, with_images=True)
    images = extract_images_from_pptx(pptx_bytes)
    assert len(images) == 1
    assert images[0] is not None


def test_extract_images_mixed():
    """PPTX with some slides having images and some without."""
    prs = Presentation()
    png_bytes = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    # Slide 1: with image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(1), Inches(1))
    # Slide 2: without image
    prs.slides.add_slide(prs.slide_layouts[6])
    # Slide 3: with image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(1), Inches(1))

    buf = io.BytesIO()
    prs.save(buf)
    images = extract_images_from_pptx(buf.getvalue())
    assert len(images) == 3
    assert images[0] is not None
    assert images[1] is None
    assert images[2] is not None


# ---------------------------------------------------------------------------
# PPTX import endpoint
# ---------------------------------------------------------------------------

def test_pptx_import_endpoint():
    """POST /api/carousel/{id}/pptx/import should extract images and update draft."""
    from fastapi import FastAPI
    from fastapi.testclient import TestClient
    from miniapp.api.routers.carousel import router
    from miniapp.api.auth import _require_auth
    from bot.services.drafts_store import DraftRecord

    app = FastAPI()
    app.include_router(router)
    app.dependency_overrides[_require_auth] = lambda: None

    draft = DraftRecord(
        draft_id="pptx01", kind="carousel", topic="Test carousel",
        source="test", created_at="2026-01-01T00:00:00",
        status="draft", feedback="",
        payload={
            "slides": ["Slide 1", "Slide 2", "Slide 3"],
            "img_prompts": ["prompt1", "prompt2", "prompt3"],
            "slide_images": [None, None, None],
        },
    )

    updated_draft = DraftRecord(
        draft_id="pptx01", kind="carousel", topic="Test carousel",
        source="test", created_at="2026-01-01T00:00:00",
        status="draft", feedback="",
        payload={"slides": ["Slide 1", "Slide 2", "Slide 3"], "images_ready": 3},
    )

    pptx_bytes = _make_test_pptx(slide_count=3, with_images=True)
    client = TestClient(app)

    with (
        patch("miniapp.api.routers.carousel.get_draft", side_effect=[draft, updated_draft]),
        patch("miniapp.api.routers.carousel.update_draft", new_callable=AsyncMock),
        patch("miniapp.api.routers.carousel.create_revision", new_callable=AsyncMock),
        patch("miniapp.api.routers.carousel.save_carousel_slide_asset", return_value={
            "filename": "slide_1_test.png",
            "url": "/generated/carousel_assets/pptx01/slide_1_test.png",
            "generated_at": "2026-01-01T00:00:00",
            "prompt": "canva_import: prompt1",
        }),
    ):
        resp = client.post(
            "/api/carousel/pptx01/pptx/import",
            files={"file": ("carousel.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )

    assert resp.status_code == 200


def test_pptx_import_empty_file():
    """Empty file should return 400."""
    from fastapi import FastAPI
    from fastapi.testclient import TestClient
    from miniapp.api.routers.carousel import router
    from miniapp.api.auth import _require_auth
    from bot.services.drafts_store import DraftRecord

    app = FastAPI()
    app.include_router(router)
    app.dependency_overrides[_require_auth] = lambda: None

    draft = DraftRecord(
        draft_id="pptx02", kind="carousel", topic="Test",
        source="test", created_at="2026-01-01T00:00:00",
        status="draft", feedback="", payload={},
    )

    client = TestClient(app)
    with patch("miniapp.api.routers.carousel.get_draft", return_value=draft):
        resp = client.post(
            "/api/carousel/pptx02/pptx/import",
            files={"file": ("empty.pptx", b"", "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert resp.status_code == 400


def test_pptx_import_not_carousel():
    """Non-carousel draft should return 404."""
    from fastapi import FastAPI
    from fastapi.testclient import TestClient
    from miniapp.api.routers.carousel import router
    from miniapp.api.auth import _require_auth
    from bot.services.drafts_store import DraftRecord

    app = FastAPI()
    app.include_router(router)
    app.dependency_overrides[_require_auth] = lambda: None

    draft = DraftRecord(
        draft_id="pptx03", kind="threads", topic="Not carousel",
        source="test", created_at="2026-01-01T00:00:00",
        status="draft", feedback="", payload={},
    )

    pptx_bytes = _make_test_pptx(slide_count=1)
    client = TestClient(app)
    with patch("miniapp.api.routers.carousel.get_draft", return_value=draft):
        resp = client.post(
            "/api/carousel/pptx03/pptx/import",
            files={"file": ("test.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert resp.status_code == 404

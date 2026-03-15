"""Tests for carousel_export_agent — PPTX export with placement + corrections."""
from __future__ import annotations

import io
import json
import tempfile
from pathlib import Path
from unittest.mock import patch, MagicMock, AsyncMock

import pytest

from bot.agents.carousel_export_agent import (
    build_pptx_from_placement,
    extract_text_positions,
    compute_corrections,
    save_corrections,
    get_aggregate_bias,
    _CORRECTIONS_LOG,
)

_EMOD = "bot.agents.carousel_export_agent"


# ── Helpers ────────────────────────────────────────────────────────────────────

def _make_test_image(w: int = 200, h: int = 200, color: str = "blue") -> bytes:
    from PIL import Image
    img = Image.new("RGB", (w, h), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ── Tests: build_pptx_from_placement ──────────────────────────────────────────

class TestBuildPptxFromPlacement:
    def test_returns_valid_pptx(self):
        slides = ["Slide 1", "Slide 2"]
        images = [_make_test_image(), _make_test_image()]
        placement = {
            "0": {"top": 0.6, "left": 0.1, "width": 0.8, "height": 0.25, "text_color": "light"},
            "1": {"top": 0.5, "left": 0.08, "width": 0.84, "height": 0.3, "text_color": "dark"},
        }

        result = build_pptx_from_placement(slides, images, placement)

        assert isinstance(result, bytes)
        assert result[:2] == b"PK"  # ZIP/PPTX magic

    def test_slide_dimensions_1080x1080(self):
        from pptx import Presentation
        from pptx.util import Emu

        slides = ["Text"]
        images = [_make_test_image()]
        placement = {"0": {"top": 0.5, "left": 0.1, "width": 0.8, "height": 0.3, "text_color": "light"}}

        pptx_bytes = build_pptx_from_placement(slides, images, placement)
        prs = Presentation(io.BytesIO(pptx_bytes))

        expected_emu = 1080 * 9525
        assert prs.slide_width == Emu(expected_emu)
        assert prs.slide_height == Emu(expected_emu)

    def test_fallback_without_placement(self):
        slides = ["No placement slide"]
        images = [_make_test_image()]

        result = build_pptx_from_placement(slides, images, placement_data=None)

        assert isinstance(result, bytes)
        assert result[:2] == b"PK"

    def test_fallback_partial_placement(self):
        slides = ["Has placement", "No placement"]
        images = [_make_test_image(), _make_test_image()]
        placement = {"0": {"top": 0.5, "left": 0.1, "width": 0.8, "height": 0.3, "text_color": "light"}}

        result = build_pptx_from_placement(slides, images, placement)

        assert result[:2] == b"PK"

    def test_no_images(self):
        slides = ["Plain text slide"]
        images = [None]

        result = build_pptx_from_placement(slides, images, placement_data=None)
        assert result[:2] == b"PK"


# ── Tests: extract_text_positions ──────────────────────────────────────────────

class TestExtractTextPositions:
    def test_extracts_positions(self):
        slides = ["Test text"]
        images = [_make_test_image()]
        placement = {"0": {"top": 0.6, "left": 0.1, "width": 0.8, "height": 0.25, "text_color": "light"}}

        pptx_bytes = build_pptx_from_placement(slides, images, placement)
        positions = extract_text_positions(pptx_bytes)

        assert isinstance(positions, list)
        assert len(positions) == 1
        assert positions[0] is not None
        for key in ("top", "left", "width", "height"):
            assert key in positions[0]

    def test_returns_none_for_empty_slide(self):
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        prs.slide_width = Emu(1080 * 9525)
        prs.slide_height = Emu(1080 * 9525)
        prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

        buf = io.BytesIO()
        prs.save(buf)

        positions = extract_text_positions(buf.getvalue())
        assert positions == [None]


# ── Tests: compute_corrections ────────────────────────────────────────────────

class TestComputeCorrections:
    def test_correct_deltas(self):
        proposed = [{"top": 0.5, "left": 0.1, "width": 0.8, "height": 0.3}]
        actual = [{"top": 0.55, "left": 0.12, "width": 0.78, "height": 0.28}]

        corrections = compute_corrections(proposed, actual)

        assert len(corrections) == 1
        c = corrections[0]
        assert c["slide_index"] == 0
        assert c["role"] == "hook"
        assert abs(c["delta"]["top"] - 0.05) < 0.001
        assert abs(c["delta"]["left"] - 0.02) < 0.001

    def test_skips_none_entries(self):
        proposed = [None, {"top": 0.5, "left": 0.1, "width": 0.8, "height": 0.3}]
        actual = [{"top": 0.6, "left": 0.1, "width": 0.8, "height": 0.3}, None]

        corrections = compute_corrections(proposed, actual)
        assert len(corrections) == 0

    def test_empty_lists(self):
        assert compute_corrections([], []) == []


# ── Tests: get_aggregate_bias ─────────────────────────────────────────────────

class TestGetAggregateBias:
    def test_returns_empty_below_threshold(self, tmp_path: Path):
        log_file = tmp_path / "corrections.jsonl"
        entries = [
            {"role": "hook", "delta": {"top": 0.05, "left": 0.02}},
            {"role": "hook", "delta": {"top": 0.03, "left": 0.01}},
        ]
        log_file.write_text("\n".join(json.dumps(e) for e in entries))

        with patch.object(
            __import__("bot.agents.carousel_export_agent", fromlist=["_CORRECTIONS_LOG"]),
            "_CORRECTIONS_LOG",
            log_file,
        ):
            from bot.agents import carousel_export_agent
            original = carousel_export_agent._CORRECTIONS_LOG
            carousel_export_agent._CORRECTIONS_LOG = log_file
            try:
                result = get_aggregate_bias("hook")
            finally:
                carousel_export_agent._CORRECTIONS_LOG = original

        assert result == {}

    def test_returns_mean_above_threshold(self, tmp_path: Path):
        log_file = tmp_path / "corrections.jsonl"
        entries = [
            {"role": "hook", "delta": {"top": 0.10, "left": 0.02, "width": 0.0, "height": 0.0}},
            {"role": "hook", "delta": {"top": 0.20, "left": 0.04, "width": 0.0, "height": 0.0}},
            {"role": "hook", "delta": {"top": 0.10, "left": 0.02, "width": 0.0, "height": 0.0}},
            {"role": "hook", "delta": {"top": 0.20, "left": 0.04, "width": 0.0, "height": 0.0}},
            {"role": "hook", "delta": {"top": 0.15, "left": 0.03, "width": 0.0, "height": 0.0}},
        ]
        log_file.write_text("\n".join(json.dumps(e) for e in entries))

        from bot.agents import carousel_export_agent
        original = carousel_export_agent._CORRECTIONS_LOG
        carousel_export_agent._CORRECTIONS_LOG = log_file
        try:
            result = get_aggregate_bias("hook")
        finally:
            carousel_export_agent._CORRECTIONS_LOG = original

        assert result != {}
        assert abs(result["top"] - 0.15) < 0.001
        assert abs(result["left"] - 0.03) < 0.001

    def test_returns_empty_for_missing_file(self, tmp_path: Path):
        from bot.agents import carousel_export_agent
        original = carousel_export_agent._CORRECTIONS_LOG
        carousel_export_agent._CORRECTIONS_LOG = tmp_path / "nonexistent.jsonl"
        try:
            result = get_aggregate_bias("hook")
        finally:
            carousel_export_agent._CORRECTIONS_LOG = original

        assert result == {}


# ── Tests: save_corrections ───────────────────────────────────────────────────

class TestSaveCorrections:
    @pytest.mark.asyncio
    async def test_appends_to_jsonl(self, tmp_path: Path):
        log_file = tmp_path / "corrections.jsonl"

        draft = MagicMock()
        draft.payload = {"slides": ["text"]}

        corrections = [
            {"slide_index": 0, "role": "hook", "delta": {"top": 0.05}},
        ]

        from bot.agents import carousel_export_agent
        original = carousel_export_agent._CORRECTIONS_LOG
        carousel_export_agent._CORRECTIONS_LOG = log_file
        try:
            with (
                patch(_EMOD + ".get_draft", new_callable=AsyncMock, return_value=draft),
                patch(_EMOD + ".update_draft", new_callable=AsyncMock),
            ):
                await save_corrections("draft123", corrections)
        finally:
            carousel_export_agent._CORRECTIONS_LOG = original

        lines = log_file.read_text().strip().split("\n")
        assert len(lines) == 1
        entry = json.loads(lines[0])
        assert entry["draft_id"] == "draft123"
        assert entry["role"] == "hook"

    @pytest.mark.asyncio
    async def test_saves_to_payload(self, tmp_path: Path):
        log_file = tmp_path / "corrections.jsonl"

        draft = MagicMock()
        draft.payload = {"slides": ["text"], "placement_corrections": []}

        corrections = [{"slide_index": 0, "role": "hook", "delta": {"top": 0.05}}]

        from bot.agents import carousel_export_agent
        original = carousel_export_agent._CORRECTIONS_LOG
        carousel_export_agent._CORRECTIONS_LOG = log_file
        try:
            with (
                patch(_EMOD + ".get_draft", new_callable=AsyncMock, return_value=draft),
                patch(_EMOD + ".update_draft", new_callable=AsyncMock) as mock_update,
            ):
                await save_corrections("draft123", corrections)
        finally:
            carousel_export_agent._CORRECTIONS_LOG = original

        call_kwargs = mock_update.call_args
        saved_payload = call_kwargs.kwargs.get("payload") or call_kwargs[1].get("payload")
        assert len(saved_payload["placement_corrections"]) == 1

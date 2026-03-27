"""Carousel Export Agent — Enhanced PPTX export with placement data + Canva correction learning."""
from __future__ import annotations

import io
import json
import logging
from pathlib import Path

from bot.services.drafts_store import get_draft, update_draft

logger = logging.getLogger(__name__)

_CORRECTIONS_LOG = Path(__file__).parent.parent.parent / "data" / "placement_corrections_log.jsonl"


# ── 0. Build PPTX for editorial layout (text as editable layer) ──────────────

def build_editorial_pptx(
    slides: list[str],
    images: list[bytes | None],
    accent_color: tuple[int, int, int] = (138, 92, 246),
    bg_color: tuple[int, int, int] = (18, 18, 22),
    username: str = "",
    avatar_bytes: bytes | None = None,
) -> bytes:
    """Build PPTX with editorial layout: photo top 55%, gradient transition, editable text.

    Uses raw (clean) images — text is NOT burnt into the image.
    Text is placed as an editable text box in the bottom section.
    Gradient between photo and text is a separate shape (editable in Canva).
    Avatar + username bar under the photo. "СВАЙПАЙ >>" on hook slide.
    """
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree

    from bot.handlers.carousel import (
        _embed_font_in_pptx,
        _FONT_NAME,
    )

    # Slide dimensions: 1080×1350 (4:5 ratio)
    SLIDE_W = 1080 * 9525  # EMU
    SLIDE_H = 1350 * 9525  # EMU
    PHOTO_FRAC = 0.55
    photo_h_emu = int(SLIDE_H * PHOTO_FRAC)
    GRADIENT_H_FRAC = 0.09  # gradient overlay height (~120px at 1350)

    BG_RGB = RGBColor(*bg_color)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = RGBColor(*accent_color)
    MUTED = RGBColor(200, 200, 210)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]

    for i, raw_text in enumerate(slides):
        text = raw_text
        if isinstance(text, list):
            text = "\n".join(str(item) for item in text)
        elif not isinstance(text, str):
            text = str(text)

        slide = prs.slides.add_slide(blank)
        is_hook = i == 0

        # Solid background for entire slide
        bg_shape = slide.shapes.add_shape(
            1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_RGB
        bg_shape.line.fill.background()

        # Photo at top — scale to full slide width, preserve aspect ratio.
        # Image may extend beyond photo_h; user can move/resize in Canva.
        img_bytes = images[i] if i < len(images) else None
        if img_bytes:
            try:
                img = Image.open(io.BytesIO(img_bytes))
                iw, ih = img.size
                # Scale to full slide width, keep aspect ratio
                scale = 1080 / iw
                new_w = 1080
                new_h = int(ih * scale)
                pic_w_emu = Emu(SLIDE_W)
                pic_h_emu = Emu(new_h * 9525)
            except Exception:
                pic_w_emu = Emu(SLIDE_W)
                pic_h_emu = Emu(photo_h_emu)

            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Emu(0), Emu(0),
                pic_w_emu, pic_h_emu,
            )

            # ── Gradient overlay (photo → bg transition) — editable shape ──
            grad_h = int(SLIDE_H * GRADIENT_H_FRAC)
            grad_top = photo_h_emu - grad_h
            grad_shape = slide.shapes.add_shape(
                1, Emu(0), Emu(grad_top), Emu(SLIDE_W), Emu(grad_h),
            )
            grad_shape.line.fill.background()
            # Build gradient fill XML: transparent at top → bg_color at bottom
            sp_pr = grad_shape._element.spPr
            # Remove any default fill
            for old_fill in sp_pr.findall(_qn("a:solidFill")) + sp_pr.findall(_qn("a:noFill")):
                sp_pr.remove(old_fill)
            grad_fill = _etree.SubElement(sp_pr, _qn("a:gradFill"))
            grad_fill.set("flip", "none")
            grad_fill.set("rotWithShape", "1")
            gs_lst = _etree.SubElement(grad_fill, _qn("a:gsLst"))
            # Stop 1: fully transparent bg_color at top (pos=0)
            gs1 = _etree.SubElement(gs_lst, _qn("a:gs"))
            gs1.set("pos", "0")
            clr1 = _etree.SubElement(gs1, _qn("a:srgbClr"))
            clr1.set("val", f"{bg_color[0]:02X}{bg_color[1]:02X}{bg_color[2]:02X}")
            alpha1 = _etree.SubElement(clr1, _qn("a:alpha"))
            alpha1.set("val", "0")
            # Stop 2: fully opaque bg_color at bottom (pos=100000)
            gs2 = _etree.SubElement(gs_lst, _qn("a:gs"))
            gs2.set("pos", "100000")
            clr2 = _etree.SubElement(gs2, _qn("a:srgbClr"))
            clr2.set("val", f"{bg_color[0]:02X}{bg_color[1]:02X}{bg_color[2]:02X}")
            alpha2 = _etree.SubElement(clr2, _qn("a:alpha"))
            alpha2.set("val", "100000")
            # Linear gradient top→bottom
            lin = _etree.SubElement(grad_fill, _qn("a:lin"))
            lin.set("ang", "5400000")  # 90 degrees = top to bottom
            lin.set("scaled", "1")

        # ── Avatar + username bar ──
        pad_x = Emu(int(SLIDE_W * 0.044))  # ~48px at 1080
        bar_top = Emu(photo_h_emu + int(SLIDE_H * 0.01))
        avatar_size_px = 56
        avatar_size_emu = avatar_size_px * 9525

        if username:
            avatar_x = pad_x
            avatar_y = bar_top

            if avatar_bytes:
                # Circular avatar from image
                try:
                    av = Image.open(io.BytesIO(avatar_bytes)).convert("RGBA")
                    av = av.resize((avatar_size_px, avatar_size_px), Image.LANCZOS)
                    # Create circular mask
                    from PIL import ImageDraw as _ImageDraw
                    mask = Image.new("L", (avatar_size_px, avatar_size_px), 0)
                    md = _ImageDraw.Draw(mask)
                    md.ellipse((0, 0, avatar_size_px, avatar_size_px), fill=255)
                    # Apply mask
                    output = Image.new("RGBA", (avatar_size_px, avatar_size_px), (0, 0, 0, 0))
                    output.paste(av, (0, 0), mask)
                    abuf = io.BytesIO()
                    output.save(abuf, format="PNG")
                    slide.shapes.add_picture(
                        io.BytesIO(abuf.getvalue()),
                        avatar_x, avatar_y,
                        Emu(avatar_size_emu), Emu(avatar_size_emu),
                    )
                except Exception:
                    avatar_bytes = None  # fall through to placeholder

            if not avatar_bytes:
                # Placeholder circle with initial letter
                circle = slide.shapes.add_shape(
                    9,  # MSO_SHAPE.OVAL
                    avatar_x, avatar_y,
                    Emu(avatar_size_emu), Emu(avatar_size_emu),
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = ACCENT
                circle.line.fill.background()
                # Add initial letter inside
                ctf = circle.text_frame
                ctf.word_wrap = False
                cp = ctf.paragraphs[0]
                cp.alignment = PP_ALIGN.CENTER
                cr = cp.add_run()
                cr.text = username[0].upper()
                cr.font.name = _FONT_NAME
                cr.font.size = Pt(20)
                cr.font.bold = True
                cr.font.color.rgb = WHITE

            # Dash line + username text
            name_x = Emu(int(pad_x) + avatar_size_emu + int(SLIDE_W * 0.013))
            # Dash
            dash_box = slide.shapes.add_textbox(
                name_x, avatar_y, Emu(int(SLIDE_W * 0.04)), Emu(avatar_size_emu),
            )
            dash_box.fill.background()
            dtf = dash_box.text_frame
            dtf.word_wrap = False
            dp = dtf.paragraphs[0]
            dp.alignment = PP_ALIGN.CENTER
            dr = dp.add_run()
            dr.text = "—"
            dr.font.name = _FONT_NAME
            dr.font.size = Pt(16)
            dr.font.color.rgb = RGBColor(100, 100, 110)

            # Username text
            uname_x = Emu(int(name_x) + int(SLIDE_W * 0.045))
            ubox = slide.shapes.add_textbox(
                uname_x, avatar_y, Emu(SLIDE_W - int(uname_x) - int(pad_x)),
                Emu(avatar_size_emu),
            )
            ubox.fill.background()
            utf = ubox.text_frame
            utf.word_wrap = False
            up = utf.paragraphs[0]
            up.alignment = PP_ALIGN.LEFT
            ur = up.add_run()
            ur.text = username
            ur.font.name = _FONT_NAME
            ur.font.size = Pt(14)
            ur.font.color.rgb = MUTED

            text_top = Emu(photo_h_emu + int(SLIDE_H * 0.055))
            text_h = Emu(SLIDE_H - photo_h_emu - int(SLIDE_H * 0.08))
        else:
            text_top = Emu(photo_h_emu + int(SLIDE_H * 0.02))
            text_h = Emu(SLIDE_H - photo_h_emu - int(SLIDE_H * 0.05))

        text_w = Emu(SLIDE_W) - pad_x * 2

        # ── Main text box (editable!) ──
        is_bullet = "\n•" in text or "\n-" in text or "\n*" in text
        display_text = text if is_bullet else text.upper()

        # Reserve space for swipe indicator
        swipe_reserve = int(SLIDE_H * 0.05) if is_hook else 0
        actual_text_h = Emu(int(text_h) - swipe_reserve)

        txBox = slide.shapes.add_textbox(pad_x, text_top, text_w, actual_text_h)
        txBox.fill.background()
        tf = txBox.text_frame
        tf.word_wrap = True

        # Determine font size
        char_count = len(text)
        if is_bullet:
            font_size = 22
        elif char_count <= 40:
            font_size = 44
        elif char_count <= 80:
            font_size = 36
        elif char_count <= 150:
            font_size = 30
        else:
            font_size = 24

        # Parse text into highlighted segments for accent coloring
        from bot.agents.carousel_editorial import _parse_highlighted_text
        uppercased = not is_bullet
        segments = _parse_highlighted_text(display_text, skip_caps_highlight=uppercased)

        p_txt = tf.paragraphs[0]
        p_txt.alignment = PP_ALIGN.CENTER if not is_bullet else PP_ALIGN.LEFT
        for seg_text, is_hl in segments:
            r = p_txt.add_run()
            r.text = seg_text
            r.font.name = _FONT_NAME
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = ACCENT if is_hl else WHITE

        # ── "СВАЙПАЙ >>" on hook slide ──
        if is_hook:
            swipe_top = Emu(SLIDE_H - int(SLIDE_H * 0.045))
            swipe_box = slide.shapes.add_textbox(
                Emu(0), swipe_top, Emu(SLIDE_W), Emu(int(SLIDE_H * 0.035)),
            )
            swipe_box.fill.background()
            stf = swipe_box.text_frame
            stf.word_wrap = False
            sp = stf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            sr = sp.add_run()
            sr.text = "СВАЙПАЙ  >>"
            sr.font.name = _FONT_NAME
            sr.font.size = Pt(18)
            sr.font.bold = False
            # Dimmed accent color
            sr.font.color.rgb = RGBColor(
                min(255, accent_color[0] // 2 + 70),
                min(255, accent_color[1] // 2 + 70),
                min(255, accent_color[2] // 2 + 70),
            )

    out = io.BytesIO()
    prs.save(out)
    return _embed_font_in_pptx(out.getvalue())


# ── 1. Build PPTX from placement data ─────────────────────────────────────────

def build_pptx_from_placement(
    slides: list[str],
    images: list[bytes | None],
    placement_data: dict[str, dict] | None = None,
) -> bytes:
    """Build PPTX using Vision-based placement_data instead of heuristic.

    Falls back to _find_text_zone() for slides without placement data.
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree

    from bot.handlers.carousel import (
        _find_text_zone,
        _embed_font_in_pptx,
        _SLIDE_EMU,
        _FONT_NAME,
    )

    BEIGE = RGBColor(0xF2, 0xE8, 0xD9)
    DARK = RGBColor(0x3D, 0x2B, 0x1F)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_EMU)
    prs.slide_height = Emu(_SLIDE_EMU)
    blank = prs.slide_layouts[6]

    for i, text in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        img_bytes = images[i] if i < len(images) else None

        # Determine placement
        pd = placement_data.get(str(i)) if placement_data else None

        if img_bytes:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes), Emu(0), Emu(0), Emu(_SLIDE_EMU), Emu(_SLIDE_EMU),
            )

            if pd:
                top_frac = pd.get("top", 0.63)
                h_frac = pd.get("height", 0.32)
                left_frac = pd.get("left", 0.08)
                w_frac = pd.get("width", 0.84)
                text_color = WHITE if pd.get("text_color", "light") == "light" else DARK
            else:
                top_frac, h_frac = _find_text_zone(img_bytes)
                left_frac, w_frac = 0.0, 1.0  # full width (matches original _build_pptx)
                text_color = WHITE
        else:
            bg = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(_SLIDE_EMU), Emu(_SLIDE_EMU))
            bg.fill.solid()
            bg.fill.fore_color.rgb = BEIGE
            bg.line.color.rgb = BEIGE
            text_color = DARK
            top_frac, h_frac = 0.32, 0.36
            left_frac, w_frac = 0.0, 1.0
            pd = None

        # Calculate positions — use placement_data left/width if available
        if pd and img_bytes:
            margin_left = Emu(int(_SLIDE_EMU * left_frac))
            box_w = Emu(int(_SLIDE_EMU * w_frac))
        else:
            margin_left = Emu(80000)
            box_w = Emu(_SLIDE_EMU) - margin_left * 2

        pad = Emu(55000)
        box_top = Emu(int(_SLIDE_EMU * top_frac))
        box_h = Emu(int(_SLIDE_EMU * h_frac))

        # Semi-transparent overlay behind text
        if img_bytes:
            overlay = slide.shapes.add_shape(
                1,
                margin_left - pad, box_top - pad,
                box_w + pad * 2, box_h + pad * 2,
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x18, 0x0E, 0x08)
            overlay.line.fill.background()
            sp_pr = overlay._element.spPr
            solid = sp_pr.find(".//" + _qn("a:solidFill"))
            if solid is not None:
                clr = solid.find(_qn("a:srgbClr"))
                if clr is not None:
                    alpha_el = _etree.SubElement(clr, _qn("a:alpha"))
                    alpha_el.set("val", "58000")

        txBox = slide.shapes.add_textbox(margin_left, box_top, box_w, box_h)
        txBox.fill.background()
        tf = txBox.text_frame
        tf.word_wrap = True

        p_txt = tf.paragraphs[0]
        align_val = pd.get("text_align", "left") if pd else "left"
        p_txt.alignment = PP_ALIGN.CENTER if align_val == "center" else PP_ALIGN.LEFT
        r_txt = p_txt.add_run()
        r_txt.text = text
        r_txt.font.name = _FONT_NAME
        r_txt.font.size = Pt(24)
        r_txt.font.bold = True
        r_txt.font.color.rgb = text_color

    out = io.BytesIO()
    prs.save(out)
    return _embed_font_in_pptx(out.getvalue())


# ── 2. Extract text positions from PPTX ───────────────────────────────────────

def extract_text_positions(pptx_bytes: bytes) -> list[dict | None]:
    """Extract text box positions from each slide as fractions of slide size.

    Returns list of dicts with top/left/width/height as fractions, or None.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_w = prs.slide_width or Emu(1080 * 9525)
    slide_h = prs.slide_height or Emu(1080 * 9525)

    results: list[dict | None] = []
    for slide in prs.slides:
        text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
        if not text_shapes:
            results.append(None)
            continue
        # Pick the shape with the most text
        shape = max(text_shapes, key=lambda s: len(s.text))
        results.append({
            "top": shape.top / slide_h,
            "left": shape.left / slide_w,
            "width": shape.width / slide_w,
            "height": shape.height / slide_h,
        })
    return results


# ── 3. Compute corrections ────────────────────────────────────────────────────

def compute_corrections(
    proposed: list[dict | None],
    actual: list[dict | None],
) -> list[dict]:
    """Compute deltas between proposed and actual placement for each slide."""
    from bot.handlers.carousel.generation import get_slide_roles

    total = max(len(proposed), len(actual))
    _roles = get_slide_roles(total)
    corrections: list[dict] = []
    for i in range(min(len(proposed), len(actual))):
        p = proposed[i]
        a = actual[i]
        if p is None or a is None:
            continue
        delta = {}
        for key in ("top", "left", "width", "height"):
            if key in p and key in a:
                delta[key] = round(a[key] - p[key], 4)
        role = _roles[i] if i < len(_roles) else "unknown"
        corrections.append({"slide_index": i, "role": role, "delta": delta})
    return corrections


# ── 4. Save corrections ──────────────────────────────────────────────────────

async def save_corrections(draft_id: str, corrections: list[dict]) -> None:
    """Save placement corrections to draft payload and JSONL log."""
    # Append to JSONL file
    _CORRECTIONS_LOG.parent.mkdir(parents=True, exist_ok=True)
    with open(_CORRECTIONS_LOG, "a") as f:
        for c in corrections:
            entry = {"draft_id": draft_id, **c}
            f.write(json.dumps(entry, ensure_ascii=False) + "\n")

    # Save to draft payload
    draft = await get_draft(draft_id)
    if not draft:
        return
    payload = dict(draft.payload)
    existing = payload.get("placement_corrections", [])
    existing.extend(corrections)
    payload["placement_corrections"] = existing
    await update_draft(draft_id, payload=payload)


# ── 5. Get aggregate bias ────────────────────────────────────────────────────

def get_aggregate_bias(role: str) -> dict:
    """Read JSONL corrections, filter by role, return mean delta.

    Returns empty dict if fewer than 5 records for this role.
    """
    if not _CORRECTIONS_LOG.exists():
        return {}

    deltas: list[dict] = []
    try:
        with open(_CORRECTIONS_LOG) as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                entry = json.loads(line)
                if entry.get("role") == role and "delta" in entry:
                    deltas.append(entry["delta"])
    except Exception:
        return {}

    if len(deltas) < 5:
        return {}

    # Compute mean delta
    keys = ("top", "left", "width", "height")
    result: dict[str, float] = {}
    for key in keys:
        values = [d[key] for d in deltas if key in d]
        if values:
            result[key] = round(sum(values) / len(values), 4)
    return result

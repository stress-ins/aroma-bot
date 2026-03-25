"""Carousel slide image generation and storage for the Mini App.

Images are saved to disk and served via /generated/carousel_assets/<draft_id>/<file>.
The draft payload is updated with slide_images: [{url, generated_at, prompt}].
"""
from __future__ import annotations

import asyncio
import functools
import logging
import os
from datetime import datetime, timezone
from pathlib import Path
from uuid import uuid4

from config import settings
from bot.services.drafts_store import get_draft, update_draft
from bot.services.gemini_images import ImageGenResult, generate_gemini_image_sync


def _get_carousel_model() -> str:
    """Get configured carousel image model, with fallback."""
    try:
        from bot.services.brand_settings_store import get_brand_settings_cached
        bs = get_brand_settings_cached()
        return bs.image_model_carousel or "gpt-image/1.5-text-to-image"
    except Exception:
        return "gpt-image/1.5-text-to-image"


def _get_img2img_model() -> str:
    """Get configured img2img model, with fallback."""
    try:
        from bot.services.brand_settings_store import get_brand_settings_cached
        bs = get_brand_settings_cached()
        return bs.image_model_img2img or "google/nano-banana-edit"
    except Exception:
        return "google/nano-banana-edit"

logger = logging.getLogger(__name__)

CAROUSEL_ASSETS_DIR = Path(
    os.getenv(
        "AROMA_CAROUSEL_ASSETS_DIR",
        Path(__file__).parent.parent.parent / "data" / "carousel_assets",
    )
)
CAROUSEL_ASSETS_DIR.mkdir(parents=True, exist_ok=True)

_FALLBACK_PROMPT = (
    "terracotta minimal lifestyle, dried herbs on warm surface, "
    "soft natural light, large negative space for text"
)


def extract_images_from_pptx(pptx_bytes: bytes) -> list[bytes | None]:
    """Extract the largest image from each slide in a PPTX file."""
    import io
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(pptx_bytes))
    images: list[bytes | None] = []
    for slide in prs.slides:
        pictures = [
            shape for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        if not pictures:
            images.append(None)
            continue
        biggest = max(pictures, key=lambda s: s.width * s.height)
        images.append(biggest.image.blob)
    return images


def _prompt_with_note(prompt: str, note: str = "") -> str:
    base = (prompt or _FALLBACK_PROMPT).strip()
    note = note.strip()
    if not note:
        return base
    return f"{base}\n\nRevision note: {note}"


def _slide_dir(draft_id: str) -> Path:
    path = CAROUSEL_ASSETS_DIR / draft_id
    path.mkdir(parents=True, exist_ok=True)
    return path


def _ensure_slide_versions(
    payload: dict[str, object],
    slide_count: int,
) -> tuple[list[dict | None], list[list[dict[str, str]]]]:
    slide_images: list[dict | None] = list(payload.get("slide_images", []))
    while len(slide_images) < slide_count:
        slide_images.append(None)

    raw_versions = payload.get("slide_image_versions", [])
    versions: list[list[dict[str, str]]] = []
    if isinstance(raw_versions, list):
        for item in raw_versions[:slide_count]:
            if isinstance(item, list):
                versions.append([
                    version
                    for version in item
                    if isinstance(version, dict) and version.get("url") and version.get("filename")
                ])
            else:
                versions.append([])
    while len(versions) < slide_count:
        versions.append([])

    for index in range(slide_count):
        current = slide_images[index]
        current_filename = str(current.get("filename", "")).strip() if isinstance(current, dict) else ""
        if current_filename and not any(str(item.get("filename", "")).strip() == current_filename for item in versions[index]):
            versions[index].append(current)

    return slide_images, versions


def save_carousel_slide_asset(
    draft_id: str,
    slide_index: int,
    image_bytes: bytes,
    *,
    prompt: str,
) -> dict[str, str]:
    directory = _slide_dir(draft_id)
    generated_at = datetime.now(timezone.utc).isoformat()
    filename = f"slide_{slide_index + 1}_{uuid4().hex[:8]}.png"
    (directory / filename).write_bytes(image_bytes)
    return {
        "filename": filename,
        "generated_at": generated_at,
        "url": f"/generated/carousel_assets/{draft_id}/{filename}",
        "prompt": prompt,
    }



def _get_blend_mood_from_payload(payload: dict) -> str | None:
    """Extract blend visual directive from draft payload if blend_context exists."""
    bc = payload.get("blend_context")
    if not bc or not isinstance(bc, dict):
        return None
    from bot.agents.blend_content_context import mood_to_visual_directive
    directive = mood_to_visual_directive(bc.get("profile", {}))
    return directive or None


async def populate_carousel_slide_assets(draft_id: str, layout_style: str = "overlay") -> str:
    """Generate Gemini images for all slides that don't have one yet.

    Updates draft payload in-place: adds/fills slide_images list.
    Runs synchronous Gemini calls sequentially (rate-limit safe).

    Returns: "done" | "awaiting_callback" | "error"
    """
    draft = await get_draft(draft_id)
    if not draft or draft.kind != "carousel":
        return "error"

    # Use layout_style from payload if not passed explicitly
    effective_layout = layout_style or draft.payload.get("layout_style", "overlay")

    img_prompts: list[str] = draft.payload.get("img_prompts", [])
    slide_images, slide_versions = _ensure_slide_versions(draft.payload, len(img_prompts))

    has_ready = False
    has_pending_callback = False
    for i, prompt in enumerate(img_prompts):
        if slide_images[i] and not isinstance(slide_images[i], dict) or (
            isinstance(slide_images[i], dict) and slide_images[i].get("filename")
        ):
            has_ready = True
            continue  # already has a real image
        if isinstance(slide_images[i], dict) and slide_images[i].get("pending_callback"):
            has_pending_callback = True
            continue  # already submitted, waiting for webhook
        try:
            result = await asyncio.get_running_loop().run_in_executor(
                None,
                functools.partial(
                    generate_gemini_image_sync,
                    prompt or _FALLBACK_PROMPT,
                    aspect_ratio="4:5",
                    log_context=f"carousel slide {i + 1}/{len(img_prompts)}",
                    model=_get_carousel_model(),
                    draft_id=draft_id,
                    content_type="carousel_slide",
                    slot_key=str(i),
                ),
            )
            if result.image_bytes:
                image_bytes = result.image_bytes
                raw_filename = None
                if effective_layout == "editorial":
                    # Save raw image (without text) for Canva export
                    raw_version = save_carousel_slide_asset(draft_id, i, image_bytes, prompt=prompt)
                    raw_filename = raw_version.get("filename")
                    slide_text = draft.payload.get("slides", [])[i] if i < len(draft.payload.get("slides", [])) else ""
                    accent = draft.payload.get("accent_color", [138, 92, 246])
                    accent_rgb = tuple(accent) if isinstance(accent, list) and len(accent) == 3 else (138, 92, 246)
                    brand_user = draft.payload.get("brand_username", "")
                    is_hook = i == 0
                    is_bullet = "\n•" in slide_text or "\n-" in slide_text or "\n*" in slide_text
                    try:
                        from bot.agents.carousel_editorial import render_editorial_png
                        image_bytes = render_editorial_png(
                            image_bytes,
                            slide_text,
                            slide_index=i,
                            accent_color=accent_rgb,
                            username=brand_user,
                            is_hook=is_hook,
                            is_bullet_list=is_bullet,
                        )
                    except Exception:
                        logger.exception("carousel_assets: editorial render failed for slide %d, using raw", i + 1)
                version = save_carousel_slide_asset(draft_id, i, image_bytes, prompt=prompt)
                if raw_filename:
                    version["raw_filename"] = raw_filename
                slide_images[i] = version
                slide_versions[i].append(version)
                has_ready = True
                logger.info("carousel_assets: slide %d generated for draft %s", i + 1, draft_id)
            elif result.kie_task_id and not result.error:
                # Submitted to Kie, waiting for webhook callback
                slide_images[i] = {"kie_task_id": result.kie_task_id, "pending_callback": True, "prompt": prompt}
                has_pending_callback = True
                logger.info("carousel_assets: slide %d submitted (awaiting callback) for draft %s", i + 1, draft_id)
            else:
                logger.warning("carousel_assets: slide %d failed for draft %s: %s", i + 1, draft_id, result.error)
        except Exception:
            logger.exception("carousel_assets: failed on slide %d for draft %s", i + 1, draft_id)

    # Always persist slide_images so the auto-trigger guard won't re-fire
    payload = dict(draft.payload)
    payload["slide_images"] = slide_images
    payload["slide_image_versions"] = slide_versions
    payload["images_ready"] = sum(
        1 for img in slide_images
        if isinstance(img, dict) and img.get("filename")
    )

    if has_pending_callback:
        payload["generation_stage"] = "awaiting_callback"
        await update_draft(draft_id, payload=payload)
        return "awaiting_callback"

    all_missing = all(
        not (isinstance(img, dict) and img.get("filename"))
        for img in slide_images
    )
    status = "error" if all_missing and not has_ready else "done"
    await update_draft(draft_id, payload=payload)
    return status


async def regenerate_carousel_slide_asset(
    draft_id: str,
    slide_index: int,
    *,
    note: str | None = None,
) -> dict[str, object] | None:
    draft = await get_draft(draft_id)
    if not draft or draft.kind != "carousel":
        return None

    img_prompts: list[str] = list(draft.payload.get("img_prompts", []))
    if slide_index < 0 or slide_index >= len(img_prompts):
        return None

    slide_images, slide_versions = _ensure_slide_versions(draft.payload, len(img_prompts))

    notes: list[str] = list(draft.payload.get("img_prompt_notes", []))
    while len(notes) < len(img_prompts):
        notes.append("")
    if note is not None:
        notes[slide_index] = note.strip()

    from bot.agents.image_prompt_router import optimize_image_prompt
    blend_mood = _get_blend_mood_from_payload(draft.payload)

    # Use img2img only when the user provided a revision note;
    # plain "regenerate" always creates a fresh image via text-to-image.
    use_img2img = bool(note and note.strip())

    image_urls: list[str] | None = None
    if use_img2img:
        current_image = slide_images[slide_index]
        if current_image and isinstance(current_image, dict):
            relative_url = current_image.get("url", "")
            base_url = settings.assets_base_url
            if relative_url and base_url:
                image_urls = [f"{base_url}{relative_url}"]

    regen_model = _get_img2img_model() if image_urls else _get_carousel_model()
    final_prompt = optimize_image_prompt(
        img_prompts[slide_index],
        model=regen_model,
        topic=draft.topic,
        slide_number=slide_index,
        total_slides=len(img_prompts),
        user_note=notes[slide_index],
        blend_mood=blend_mood,
    )

    try:
        result = await asyncio.get_running_loop().run_in_executor(
            None,
            functools.partial(
                generate_gemini_image_sync,
                final_prompt,
                aspect_ratio="4:5",
                image_urls=image_urls,
                log_context=f"carousel slide regenerate {slide_index + 1}/{len(img_prompts)}",
                model=regen_model,
                draft_id=draft_id,
                content_type="carousel_slide",
                slot_key=str(slide_index),
            ),
        )
    except Exception:
        logger.exception("carousel_assets: regenerate failed on slide %d for draft %s", slide_index + 1, draft_id)
        return None

    if not result.image_bytes:
        return None

    version = save_carousel_slide_asset(draft_id, slide_index, result.image_bytes, prompt=final_prompt)
    slide_images[slide_index] = version
    slide_versions[slide_index].append(version)
    payload = dict(draft.payload)
    payload["slide_images"] = slide_images
    payload["slide_image_versions"] = slide_versions
    payload["img_prompt_notes"] = notes
    payload["images_ready"] = sum(1 for img in slide_images if img)
    updated = await update_draft(draft_id, payload=payload)
    return dict(updated.payload) if updated else None


async def regenerate_all_carousel_slide_assets(draft_id: str) -> dict[str, object] | None:
    draft = await get_draft(draft_id)
    if not draft or draft.kind != "carousel":
        return None

    img_prompts: list[str] = list(draft.payload.get("img_prompts", []))
    slide_images, slide_versions = _ensure_slide_versions(draft.payload, len(img_prompts))
    notes: list[str] = list(draft.payload.get("img_prompt_notes", []))
    while len(notes) < len(img_prompts):
        notes.append("")

    from bot.agents.image_prompt_router import optimize_image_prompt as _optimize
    blend_mood = _get_blend_mood_from_payload(draft.payload)

    changed = False
    for index, prompt in enumerate(img_prompts):
        # regenerate-all always uses text-to-image (fresh generation)
        image_urls = None
        regen_model = _get_carousel_model()
        try:
            final_prompt = _optimize(
                prompt, model=regen_model, topic=draft.topic,
                slide_number=index, total_slides=len(img_prompts),
                user_note=notes[index], blend_mood=blend_mood,
            )
        except Exception:
            logger.exception("carousel_assets: optimize failed on slide %d", index + 1)
            final_prompt = _prompt_with_note(prompt, notes[index])

        try:
            result = await asyncio.get_running_loop().run_in_executor(
                None,
                functools.partial(
                    generate_gemini_image_sync,
                    final_prompt,
                    aspect_ratio="4:5",
                    image_urls=image_urls,
                    log_context=f"carousel slide regenerate all {index + 1}/{len(img_prompts)}",
                    model=regen_model,
                    draft_id=draft_id,
                    content_type="carousel_slide",
                    slot_key=str(index),
                ),
            )
        except Exception:
            logger.exception("carousel_assets: regenerate-all failed on slide %d for draft %s", index + 1, draft_id)
            continue
        if not result.image_bytes:
            continue
        version = save_carousel_slide_asset(draft_id, index, result.image_bytes, prompt=final_prompt)
        slide_images[index] = version
        slide_versions[index].append(version)
        changed = True

    if not changed:
        return dict(draft.payload)

    payload = dict(draft.payload)
    payload["slide_images"] = slide_images
    payload["slide_image_versions"] = slide_versions
    payload["img_prompt_notes"] = notes
    payload["images_ready"] = sum(1 for img in slide_images if img)
    updated = await update_draft(draft_id, payload=payload)
    return dict(updated.payload) if updated else None


async def select_carousel_slide_version(
    draft_id: str,
    slide_index: int,
    version_index: int,
) -> dict[str, object] | None:
    draft = await get_draft(draft_id)
    if not draft or draft.kind != "carousel":
        return None

    prompts: list[str] = list(draft.payload.get("img_prompts", []))
    if slide_index < 0 or slide_index >= len(prompts):
        return None

    slide_images, slide_versions = _ensure_slide_versions(draft.payload, len(prompts))
    if version_index < 0 or version_index >= len(slide_versions[slide_index]):
        return None

    slide_images[slide_index] = slide_versions[slide_index][version_index]
    payload = dict(draft.payload)
    payload["slide_images"] = slide_images
    payload["slide_image_versions"] = slide_versions
    payload["images_ready"] = sum(1 for img in slide_images if img)
    updated = await update_draft(draft_id, payload=payload)
    return dict(updated.payload) if updated else None


async def delete_carousel_slide_version(
    draft_id: str,
    slide_index: int,
    version_index: int,
) -> dict[str, object] | None:
    draft = await get_draft(draft_id)
    if not draft or draft.kind != "carousel":
        return None

    prompts: list[str] = list(draft.payload.get("img_prompts", []))
    if slide_index < 0 or slide_index >= len(prompts):
        return None

    slide_images, slide_versions = _ensure_slide_versions(draft.payload, len(prompts))
    if version_index < 0 or version_index >= len(slide_versions[slide_index]):
        return None

    removed = slide_versions[slide_index].pop(version_index)
    current = slide_images[slide_index]
    current_filename = str(current.get("filename", "")).strip() if isinstance(current, dict) else ""
    removed_filename = str(removed.get("filename", "")).strip()
    if current_filename and current_filename == removed_filename:
      slide_images[slide_index] = slide_versions[slide_index][-1] if slide_versions[slide_index] else None

    payload = dict(draft.payload)
    payload["slide_images"] = slide_images
    payload["slide_image_versions"] = slide_versions
    payload["images_ready"] = sum(1 for img in slide_images if img)
    updated = await update_draft(draft_id, payload=payload)
    return dict(updated.payload) if updated else None


async def update_carousel_slide_text(
    draft_id: str,
    slide_index: int,
    text: str,
) -> dict[str, object] | None:
    return await _update_carousel_list_field(draft_id, "slides", slide_index, text)


async def update_carousel_slide_note(
    draft_id: str,
    slide_index: int,
    note: str,
) -> dict[str, object] | None:
    return await _update_carousel_list_field(
        draft_id, "img_prompt_notes", slide_index, note, pad_to="img_prompts"
    )


async def _update_carousel_list_field(
    draft_id: str,
    field_key: str,
    index: int,
    value: str,
    *,
    pad_to: str | None = None,
) -> dict | None:
    draft = await get_draft(draft_id)
    if not draft or draft.kind != "carousel":
        return None
    items: list = list(draft.payload.get(field_key, []))
    if pad_to:
        ref_len = len(draft.payload.get(pad_to, []))
        while len(items) < ref_len:
            items.append("")
    if index < 0 or index >= len(items):
        return None
    items[index] = value.strip()
    payload = dict(draft.payload)
    payload[field_key] = items
    updated = await update_draft(draft_id, payload=payload)
    return dict(updated.payload) if updated else None


def load_carousel_slide_images(draft_id: str, slide_images: list[dict | None]) -> list[bytes | None]:
    images: list[bytes | None] = []
    for item in slide_images:
        if not item:
            images.append(None)
            continue
        filename = str(item.get("filename", "")).strip()
        if not filename:
            images.append(None)
            continue
        path = _slide_dir(draft_id) / filename
        images.append(path.read_bytes() if path.exists() else None)
    return images


def load_carousel_raw_images(draft_id: str, slide_images: list[dict | None]) -> list[bytes | None]:
    """Load raw (without editorial text overlay) images for Canva/PPTX export.

    Falls back to the main filename if raw_filename is not available.
    """
    images: list[bytes | None] = []
    for item in slide_images:
        if not item:
            images.append(None)
            continue
        raw_fn = str(item.get("raw_filename", "")).strip()
        fn = str(item.get("filename", "")).strip()
        filename = raw_fn or fn
        if not filename:
            images.append(None)
            continue
        path = _slide_dir(draft_id) / filename
        images.append(path.read_bytes() if path.exists() else None)
    return images

"""PPTX generation for carousel slides."""
from __future__ import annotations

import io
import logging
import zipfile

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from bot.handlers.carousel.generation import (
    _FONT_PATH,
    _FONT_NAME,
    _SLIDE_EMU,
    _find_text_zone,
    wrap_slide_text,
)

logger = logging.getLogger(__name__)

_FONT_REL_ID = "rIdDeldedaRegular"


def _embed_font_in_pptx(pptx_bytes: bytes) -> bytes:
    """Embed DeldedaOpen.ttf into the PPTX ZIP so the font travels with the file."""
    if not _FONT_PATH.exists():
        logger.warning("Font file not found: %s -- skipping font embedding", _FONT_PATH)
        return pptx_bytes

    font_data = _FONT_PATH.read_bytes()
    inp = io.BytesIO(pptx_bytes)
    out = io.BytesIO()

    font_rel_entry = (
        f'<Relationship Id="{_FONT_REL_ID}" '
        f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
        f'Target="fonts/font1.ttf"/>'
    ).encode()

    font_xml_entry = (
        f'<p:embeddedFontLst>'
        f'<p:embeddedFont>'
        f'<p:font typeface="{_FONT_NAME}" charset="0" pitchFamily="32"/>'
        f'<p:regular r:id="{_FONT_REL_ID}"/>'
        f'</p:embeddedFont>'
        f'</p:embeddedFontLst>'
    ).encode()

    content_type_entry = (
        b'<Override PartName="/ppt/fonts/font1.ttf" '
        b'ContentType="application/x-fontdata"/>'
    )

    with zipfile.ZipFile(inp, "r") as zin, \
         zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:

        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename == "ppt/_rels/presentation.xml.rels":
                data = data.replace(b"</Relationships>", font_rel_entry + b"</Relationships>")

            elif item.filename == "ppt/presentation.xml":
                data = data.replace(b"</p:presentation>", font_xml_entry + b"</p:presentation>")

            elif item.filename == "[Content_Types].xml":
                data = data.replace(b"</Types>", content_type_entry + b"</Types>")

            zout.writestr(item, data)

        # Add the font binary
        zout.writestr("ppt/fonts/font1.ttf", font_data)

    return out.getvalue()


def _bake_overlay_rect(
    img_bytes: bytes,
    top_frac: float,
    h_frac: float,
    margin_frac: float = 80000 / _SLIDE_EMU,
    width_frac: float | None = None,
) -> bytes:
    """Burn a semi-transparent dark overlay rectangle into image pixels.

    This avoids a separate PPTX shape for the overlay which Canva
    consistently moves to the wrong layer on import.
    """
    img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    w, h = img.size

    pad_frac = 55000 / _SLIDE_EMU
    if width_frac is None:
        width_frac = 1.0 - 2 * margin_frac

    box_top = int(h * top_frac)
    box_h = int(h * h_frac)
    pad_x = int(w * pad_frac)
    pad_y = int(h * pad_frac)
    margin_px = int(w * margin_frac)
    box_w_px = int(w * width_frac)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    alpha = int(255 * 0.58)
    x0 = margin_px - pad_x
    y0 = box_top - pad_y
    x1 = margin_px + box_w_px + pad_x
    y1 = box_top + box_h + pad_y
    draw.rounded_rectangle([x0, y0, x1, y1], radius=12, fill=(0x18, 0x0E, 0x08, alpha))

    composite = Image.alpha_composite(img, overlay)
    buf = io.BytesIO()
    composite.convert("RGB").save(buf, format="JPEG", quality=92)
    return buf.getvalue()


def _build_pptx(slides: list[str], images: list[bytes | None] | None = None) -> bytes:
    BEIGE = RGBColor(0xF2, 0xE8, 0xD9)
    DARK  = RGBColor(0x3D, 0x2B, 0x1F)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width  = Emu(_SLIDE_EMU)
    prs.slide_height = Emu(_SLIDE_EMU)
    blank = prs.slide_layouts[6]

    for i, raw_text in enumerate(slides):
        text = raw_text
        if isinstance(text, list):
            text = "\n".join(str(item) for item in text)
        elif not isinstance(text, str):
            text = str(text)
        slide = prs.slides.add_slide(blank)
        img_bytes = (images[i] if images and i < len(images) else None)

        if img_bytes:
            top_frac, h_frac = _find_text_zone(img_bytes)
            # Bake semi-transparent overlay directly into the image so Canva
            # cannot reorder layers (overlay was a separate shape before and
            # Canva consistently moved it on top of everything).
            composite_bytes = _bake_overlay_rect(img_bytes, top_frac, h_frac)
            slide.shapes.add_picture(
                io.BytesIO(composite_bytes), Emu(0), Emu(0), Emu(_SLIDE_EMU), Emu(_SLIDE_EMU)
            )
            text_color = WHITE
        else:
            bg = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(_SLIDE_EMU), Emu(_SLIDE_EMU))
            bg.fill.solid()
            bg.fill.fore_color.rgb = BEIGE
            bg.line.color.rgb = BEIGE
            text_color = DARK
            top_frac, h_frac = 0.32, 0.36   # centre for plain background

        margin  = Emu(80000)
        box_top = Emu(int(_SLIDE_EMU * top_frac))
        box_h   = Emu(int(_SLIDE_EMU * h_frac))
        box_w   = Emu(_SLIDE_EMU) - margin * 2

        txBox = slide.shapes.add_textbox(
            margin, box_top, box_w, box_h
        )
        txBox.fill.background()
        tf = txBox.text_frame
        tf.word_wrap = True

        # Pre-wrap text using shared wrapping logic for consistency with preview
        wrapped_lines = wrap_slide_text(text, max_chars_per_line=32)
        wrapped_text = "\n".join(wrapped_lines)

        p_txt = tf.paragraphs[0]
        p_txt.alignment = PP_ALIGN.LEFT
        r_txt = p_txt.add_run()
        r_txt.text = wrapped_text
        r_txt.font.name = _FONT_NAME
        r_txt.font.size = Pt(24)
        r_txt.font.bold = True
        r_txt.font.color.rgb = text_color

    out = io.BytesIO()
    prs.save(out)
    return _embed_font_in_pptx(out.getvalue())

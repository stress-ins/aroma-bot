from __future__ import annotations

import asyncio
import html as _html
import io
import logging
import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, InputMediaPhoto
from telegram.ext import ContextTypes, CommandHandler, CallbackQueryHandler, MessageHandler, filters

from bot.services.gemini_images import generate_gemini_image_sync
from config import settings
from bot.handlers.threads import _format_trends, _claude_topics, _fix_dashes

logger = logging.getLogger(__name__)

_executor = ThreadPoolExecutor(max_workers=2)
_img_executor = ThreadPoolExecutor(max_workers=5)

# ── Font ────────────────────────────────────────────────────────────────────
_FONT_PATH = Path(__file__).parent.parent.parent / "assets" / "fonts" / "DeldedaOpen.ttf"
_FONT_NAME = "Deledda Open"
_FONT_REL_ID = "rIdDeldedaRegular"

# ── Layout constants ────────────────────────────────────────────────────────
_SLIDE_EMU = 1080 * 9525

_DEFAULT_FORBIDDEN_VISUAL_MOTIFS = [
    "hands joined together",
    "prayer pose",
    "praying hands",
    "cult ritual circle",
    "sect-like worship",
    "two cupped hands around one object",
]

# Slide narrative roles — used for regen-slide context
_SLIDE_VISUAL_ROLES = [
    "hook — powerful first image that stops scrolling",
    "problem — a relatable moment of stress, overload, or disconnection",
    "mechanism — body process shown metaphorically, a transitional mood",
    "insight — moment of clarity, first hint of the brand world",
    "solution — the sensory practice, the remedy in action",
    "call to action — warm human invitation",
]

_SLIDE_LABELS = [
    "Слайд 1 — Hook",
    "Слайд 2 — Проблема",
    "Слайд 3 — Механизм",
    "Слайд 4 — Инсайт",
    "Слайд 5 — Решение",
    "Слайд 6 — CTA",
]

# ── Prompts ─────────────────────────────────────────────────────────────────
_PROMPT_CAROUSEL = """\
Ты — сценарист карусели для Instagram. Ниша: регуляция нервной системы через \
сенсорные практики (ароматерапия, медитации, гонг).
Создай черновик карусели из 5 слайдов по теме: {topic}

Требования:
- Описывай только реалистичные повседневные ситуации. Все детали одного слайда должны быть логически связаны и физически возможны
- Слайд 1: цепляющий хук, до 60 символов
- Слайды 2-4: один тезис + 1-2 предложения, до 120 символов на слайд
- Слайд 5: призыв к действию, до 80 символов
- Живой язык, от первого лица, без клише и длинных тире
- Базовый промпт для фото (английский, 15-25 слов): сцена, палитра из темы, освещение, композиция. Без параметров Midjourney.

Формат — строго:
SLIDE1: [текст]
SLIDE2: [текст]
SLIDE3: [текст]
SLIDE4: [текст]
SLIDE5: [текст]
IMG_PROMPT: [промпт]
"""

_PROMPT_TOPICS = """\
Ты — стратег по контенту в Instagram. Ниша: регуляция нервной системы через \
сенсорные практики (ароматерапия, медитации, гонг).

На основе трендов ниже предложи 10 тем для карусели. \
Каждая тема — конкретный угол, ситуация или вопрос. Без воды и банальностей.

Формат — строго нумерованный список:
1. [тема]
...
10. [тема]
"""


# ── Claude helpers ──────────────────────────────────────────────────────────

def _claude_topics_carousel(trends_text: str) -> list[str]:
    from bot.services.claude_client import call_claude
    text = call_claude(
        messages=[{"role": "user", "content": f"Тренды:\n{trends_text}"}],
        max_tokens=800,
        system=_PROMPT_TOPICS,
        context="carousel topics",
    )
    topics: list[str] = []
    for line in text.splitlines():
        line = line.strip()
        if line and line[0].isdigit() and ". " in line:
            topics.append(line.split(". ", 1)[1].strip())
    return topics[:10]


def _claude_carousel_draft(topic: str, angle: str = "", hook: str = "") -> tuple[list[str], str]:
    from bot.agents.carousel_editor import _parse_slides
    from bot.services.claude_client import call_claude
    user_content = _PROMPT_CAROUSEL.format(topic=topic)
    if angle or hook:
        user_content += f"\n\nСтратег предложил:\nУгол: {angle}\nХук: {hook}\nИспользуй этот угол и хук как основу."
    text = call_claude(
        messages=[{"role": "user", "content": user_content}],
        max_tokens=900,
        context="carousel draft",
    )
    slides = [_fix_dashes(s) for s in _parse_slides(text, count=5)]
    img_prompt = ""
    for line in text.splitlines():
        line = line.strip()
        if line.upper().startswith("IMG_PROMPT:"):
            img_prompt = line.split(":", 1)[1].strip()
            break
    return slides, img_prompt



def _forbidden_visual_motifs() -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for motif in [*_DEFAULT_FORBIDDEN_VISUAL_MOTIFS, *getattr(settings, "carousel_forbidden_visual_motifs_list", [])]:
        value = str(motif).strip()
        lowered = value.lower()
        if not value or lowered in seen:
            continue
        seen.add(lowered)
        result.append(value)
    return result


def _forbidden_visual_motifs_text() -> str:
    return ", ".join(_forbidden_visual_motifs())


def _generate_slide_image_prompts_sync(slides: list[str], topic: str, blend_context: dict | None = None) -> list[str]:
    """Generate one unique, detailed image prompt per slide via art-director + image prompt router.

    Step 1: Claude as art director produces short 30-50 word prompts per slide.
    Step 2: Each prompt is expanded by the image prompt router (model-aware).
    """
    from bot.agents.image_prompt_router import optimize_image_prompt
    from bot.services.claude_client import call_claude

    slides_desc = "\n".join(
        f"Slide {i + 1}: {text}" for i, text in enumerate(slides)
    )

    blend_visual = ""
    if blend_context:
        from bot.agents.blend_content_context import mood_to_visual_directive
        visual_directive = mood_to_visual_directive(blend_context.get("profile", {}))
        oil_names_en = [o.get("name_en", "") for o in blend_context.get("oils", []) if o.get("name_en")]
        if visual_directive:
            blend_visual = f"\nBlend visual mood: {visual_directive}\n"
        if oil_names_en:
            blend_visual += f"Featured botanicals: {', '.join(oil_names_en)}\n"

    prompt = (
        f'You are an art director for an Instagram carousel on the topic: "{topic}"\n\n'
        f"{blend_visual}"
        "Step 1: Determine the ideal color palette, mood, and lighting from the topic.\n"
        "The palette must match the emotional tone of the topic. Examples:\n"
        "- 'кайф от кабриолета' → bright sun, azure sky, warm golden light, vivid saturated colors\n"
        "- 'заземление через аромат' → earth tones, terracotta, sage, warm wood\n"
        "- 'зимний уют' → deep amber, warm burgundy, soft candlelight glow\n\n"
        "Step 2: Determine each slide's narrative role (hook, problem, mechanism, insight, solution, CTA).\n\n"
        "Step 3: Generate one image prompt per slide. Each must:\n"
        "- Be 30-50 words in English\n"
        "- Describe the scene, palette (from your Step 1 analysis), lighting, and composition\n"
        "- Be visually distinct from other slides\n"
        "- NOT use any Midjourney parameters (no --ar, --style, --v, --no)\n\n"
        "Universal constraints:\n"
        "- Color intensity: rich, saturated tones matching topic mood. Never grey/washed-out.\n"
        "- Forbidden everywhere: stock photo look, plastic, harsh shadows, "
        "white/grey plain backgrounds, human faces, any text or typography\n"
        f"- Also forbidden everywhere: {_forbidden_visual_motifs_text()}\n"
        "- Required: leave a large clear area (at least 1/3 of frame) for text overlay\n\n"
        f"{slides_desc}\n\n"
        "Return strictly in this format, nothing else:\n"
        + "\n".join(f"IMG{i + 1}: [prompt]" for i in range(len(slides)))
    )

    raw_text = call_claude(
        messages=[{"role": "user", "content": prompt}],
        max_tokens=900,
        context="carousel image prompts",
    )

    parsed: dict[int, str] = {}
    for line in raw_text.splitlines():
        line = line.strip()
        for i in range(1, len(slides) + 1):
            if line.startswith(f"IMG{i}:"):
                parsed[i - 1] = line.split(":", 1)[1].strip()
                break

    raw_prompts: list[str] = []
    for i, slide in enumerate(slides):
        if i in parsed:
            raw_prompts.append(parsed[i])
        else:
            raw_prompts.append(
                f"vibrant lifestyle scene matching '{topic}', {slide[:35]}, "
                f"rich saturated colors, natural light"
            )

    # Expand each raw prompt through the image prompt router
    from bot.services.carousel_assets import _get_carousel_model
    blend_mood = ""
    if blend_context:
        from bot.agents.blend_content_context import mood_to_visual_directive as _m2v
        blend_mood = _m2v(blend_context.get("profile", {}))

    model = _get_carousel_model()
    optimized: list[str] = []
    for i, raw in enumerate(raw_prompts):
        try:
            expanded = optimize_image_prompt(
                raw, model=model, topic=topic, slide_number=i, total_slides=len(slides),
                blend_mood=blend_mood or None,
            )
            optimized.append(expanded)
        except Exception:
            logger.exception("Image prompt router failed for slide %d", i + 1)
            optimized.append(raw)
    return optimized


def _generate_carousel_sync(topic: str, user_forbidden: list[str] | None = None, blend_context: dict | None = None) -> tuple[list[str], list[str], str, str]:
    """Draft → editor → 6 refined slides + per-slide image prompts.

    Returns (slides, img_prompts, angle, hook).
    """
    from bot.agents.carousel_editor import edit_carousel_sync
    import time

    from bot.agents.content import _generate_strategist_sync
    angle, hook = _generate_strategist_sync(topic, goal_key="trust", format_key="carousel", blend_context=blend_context)

    for attempt in range(2):
        try:
            raw_slides, _ = _claude_carousel_draft(topic, angle, hook)
            if not raw_slides:
                logger.warning("_claude_carousel_draft empty on attempt %d, topic: %s", attempt + 1, topic)
                if attempt == 0:
                    time.sleep(3)
                    continue
                return [], [], angle, hook
            refined = edit_carousel_sync(raw_slides, topic, user_forbidden=user_forbidden or [])
            if not refined:
                logger.warning("edit_carousel_sync empty on attempt %d, topic: %s", attempt + 1, topic)
                if attempt == 0:
                    time.sleep(3)
                    continue
                return [], [], angle, hook
            img_prompts = _generate_slide_image_prompts_sync(refined, topic, blend_context=blend_context)
            return refined, img_prompts, angle, hook
        except Exception:
            logger.exception("_generate_carousel_sync attempt %d failed for topic: %s", attempt + 1, topic)
            if attempt == 0:
                time.sleep(3)
                continue
    return [], [], angle, hook


# ── Image analysis ──────────────────────────────────────────────────────────

def _find_text_zone(img_bytes: bytes) -> tuple[float, float]:
    """Return (top_fraction, height_fraction) for text box placement.
    Scans 5 horizontal bands; picks the calmest (low variance) + darkest one —
    low variance means no busy detail to compete with text,
    dark means white text will be readable.
    """
    try:
        from PIL import Image as _PIL, ImageStat, ImageFilter
        img = _PIL.open(io.BytesIO(img_bytes)).convert("L").resize((90, 90))
        img = img.filter(ImageFilter.GaussianBlur(2))
        W, H = img.size  # 90, 90
        n = 5
        bh = H // n
        best_score = float("inf")
        best_band = 3  # safe default: bottom-ish
        for b in range(n):
            band = img.crop((0, b * bh, W, (b + 1) * bh))
            stat = ImageStat.Stat(band)
            avg = stat.mean[0]       # 0 = black, 255 = white
            std = stat.stddev[0]     # 0 = uniform, high = busy
            # Lower score = better: prefer calm (low std) and dark (low avg)
            score = std * 1.5 + avg * 0.4
            if score < best_score:
                best_score = score
                best_band = b
        top_frac = best_band / n + 0.01
        h_frac   = 1 / n + 0.04          # slightly taller than one band
        if top_frac + h_frac > 0.97:
            top_frac = 0.97 - h_frac
        return top_frac, h_frac
    except Exception:
        return 0.63, 0.32  # safe default: lower third


def _apply_note_to_prompt(prompt: str, note: str) -> str:
    """Combine prompt and user note."""
    base = prompt.strip().rstrip(",")
    return f"{base}, {note.strip()}" if note.strip() else base


def _qa_image_sync(
    img_bytes: bytes, prompt: str, note: str = "", slide_idx: int = -1
) -> tuple[bool, str]:
    """Vision QA: check for hallucinations, forbidden elements, note compliance."""
    from bot.services.claude_client import call_claude
    import base64

    note_check = (
        f"\n4. The user requested: \"{note}\" — verify this is clearly reflected."
        if note else ""
    )
    qa_prompt = (
        f"You are a strict visual QA agent for Instagram carousel images.\n"
        f"Image prompt used: {prompt}\n\n"
        f"Check this image for issues:\n"
        f"1. Physically impossible or hallucinated elements "
        f"(e.g. lavender on fire, smoke from cold objects, impossible anatomy of objects)\n"
        f"2. Any visible text, watermarks, logos, or typography"
        f"\n2b. Any forbidden visual motifs such as: {_forbidden_visual_motifs_text()}"
        f"{note_check}\n\n"
        f"Reply in this exact format (2 lines only):\n"
        f"PASS or FAIL\n"
        f"REASON: [one short sentence. If PASS write: OK]"
    )
    try:
        text = call_claude(
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/jpeg",
                            "data": base64.standard_b64encode(img_bytes).decode(),
                        },
                    },
                    {"type": "text", "text": qa_prompt},
                ],
            }],
            max_tokens=80,
            context="carousel qa image",
        )
        passed = text.upper().startswith("PASS")
        reason = ""
        for line in text.splitlines():
            if line.upper().startswith("REASON:"):
                reason = line.split(":", 1)[1].strip()
                break
        return passed, reason
    except Exception as exc:
        logger.warning("QA vision error: %s", exc)
        return True, ""  # Don't block on QA errors


# ── Gemini ──────────────────────────────────────────────────────────────────

def _gemini_slide(prompt: str, key_index: int = 0) -> bytes | None:
    return generate_gemini_image_sync(prompt, log_context="Gemini carousel")


# ── PPTX ────────────────────────────────────────────────────────────────────

def _embed_font_in_pptx(pptx_bytes: bytes) -> bytes:
    """Embed DeldedaOpen.ttf into the PPTX ZIP so the font travels with the file."""
    if not _FONT_PATH.exists():
        logger.warning("Font file not found: %s — skipping font embedding", _FONT_PATH)
        return pptx_bytes

    font_data = _FONT_PATH.read_bytes()
    inp = io.BytesIO(pptx_bytes)
    out = io.BytesIO()

    font_rel_entry = (
        f'<Relationship Id="{_FONT_REL_ID}" '
        f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
        f'Target="fonts/font1.ttf"/>'
    ).encode()

    font_xml_entry = (
        f'<p:embeddedFontLst>'
        f'<p:embeddedFont>'
        f'<p:font typeface="{_FONT_NAME}" charset="0" pitchFamily="32"/>'
        f'<p:regular r:id="{_FONT_REL_ID}"/>'
        f'</p:embeddedFont>'
        f'</p:embeddedFontLst>'
    ).encode()

    content_type_entry = (
        b'<Override PartName="/ppt/fonts/font1.ttf" '
        b'ContentType="application/x-fontdata"/>'
    )

    with zipfile.ZipFile(inp, "r") as zin, \
         zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:

        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename == "ppt/_rels/presentation.xml.rels":
                data = data.replace(b"</Relationships>", font_rel_entry + b"</Relationships>")

            elif item.filename == "ppt/presentation.xml":
                data = data.replace(b"</p:presentation>", font_xml_entry + b"</p:presentation>")

            elif item.filename == "[Content_Types].xml":
                data = data.replace(b"</Types>", content_type_entry + b"</Types>")

            zout.writestr(item, data)

        # Add the font binary
        zout.writestr("ppt/fonts/font1.ttf", font_data)

    return out.getvalue()


def _build_pptx(slides: list[str], images: list[bytes | None] | None = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BEIGE = RGBColor(0xF2, 0xE8, 0xD9)
    DARK  = RGBColor(0x3D, 0x2B, 0x1F)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width  = Emu(_SLIDE_EMU)
    prs.slide_height = Emu(_SLIDE_EMU)
    blank = prs.slide_layouts[6]

    for i, text in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        img_bytes = (images[i] if images and i < len(images) else None)

        if img_bytes:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes), Emu(0), Emu(0), Emu(_SLIDE_EMU), Emu(_SLIDE_EMU)
            )
            text_color = WHITE
            top_frac, h_frac = _find_text_zone(img_bytes)
        else:
            bg = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(_SLIDE_EMU), Emu(_SLIDE_EMU))
            bg.fill.solid()
            bg.fill.fore_color.rgb = BEIGE
            bg.line.color.rgb = BEIGE
            text_color = DARK
            top_frac, h_frac = 0.32, 0.36   # centre for plain background

        pad     = Emu(55000)
        margin  = Emu(80000)
        box_top = Emu(int(_SLIDE_EMU * top_frac))
        box_h   = Emu(int(_SLIDE_EMU * h_frac))

        # Semi-transparent dark overlay behind text (only when image is present)
        if img_bytes:
            from pptx.oxml.ns import qn as _qn
            from lxml import etree as _etree
            overlay = slide.shapes.add_shape(
                1,
                margin - pad, box_top - pad,
                Emu(_SLIDE_EMU) - margin * 2 + pad * 2, box_h + pad * 2,
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x18, 0x0E, 0x08)
            overlay.line.fill.background()
            # Set 58% opacity via OOXML (val=58000 means 58% opaque)
            sp_pr = overlay._element.spPr
            solid = sp_pr.find(".//" + _qn("a:solidFill"))
            if solid is not None:
                clr = solid.find(_qn("a:srgbClr"))
                if clr is not None:
                    alpha_el = _etree.SubElement(clr, _qn("a:alpha"))
                    alpha_el.set("val", "58000")

        txBox = slide.shapes.add_textbox(
            margin, box_top, Emu(_SLIDE_EMU) - margin * 2, box_h
        )
        txBox.fill.background()
        tf = txBox.text_frame
        tf.word_wrap = True

        # Only the slide text — no labels
        p_txt = tf.paragraphs[0]
        p_txt.alignment = PP_ALIGN.LEFT
        r_txt = p_txt.add_run()
        r_txt.text = text
        r_txt.font.name = _FONT_NAME
        r_txt.font.size = Pt(24)
        r_txt.font.bold = True
        r_txt.font.color.rgb = text_color

    out = io.BytesIO()
    prs.save(out)
    return _embed_font_in_pptx(out.getvalue())


# ── Text formatters ──────────────────────────────────────────────────────────

def _make_slide_prompts_with_text(img_prompts: list[str], slides: list[str]) -> str:
    """HTML-formatted — send with parse_mode=HTML.
    Each slide gets its own unique image prompt with text overlay injected.
    """
    lines = ["<b>🍌 Nana Banana — с текстом (1080×1350, --ar 4:5):</b>\n"]
    flags = "--ar 4:5 --style atmospheric"
    for i, (prompt, slide) in enumerate(zip(img_prompts, slides), 1):
        # Strip trailing flags, inject text overlay, re-add flags
        clean = prompt.replace(flags, "").strip().rstrip(",").rstrip()
        full_prompt = f'{clean}, text overlay: "{slide[:50]}", {flags}'
        lines.append(
            f"<b>Слайд {i}:</b>\n"
            f"<pre>{_html.escape(full_prompt)}</pre>"
        )
    return "\n".join(lines)


def _make_slide_prompts_no_text(img_prompts: list[str], slides: list[str]) -> str:
    """HTML-formatted — send with parse_mode=HTML.
    Each slide gets its own unique image prompt (backgrounds, no text).
    """
    lines = ["<b>🍌 Nana Banana — фон без текста (1080×1350, --ar 4:5):</b>\n"]
    for i, (prompt, slide) in enumerate(zip(img_prompts, slides), 1):
        lines.append(
            f"<b>Слайд {i}</b> — {_html.escape(slide[:45])}\n"
            f"<pre>{_html.escape(prompt)}</pre>"
        )
    return "\n".join(lines)


def _format_for_canva(slides: list[str]) -> str:
    """HTML-formatted — send with parse_mode=HTML."""
    parts = ["<b>📋 Тексты для Canva:</b>\n"]
    for i, slide in enumerate(slides):
        label = _SLIDE_LABELS[i] if i < len(_SLIDE_LABELS) else f"Слайд {i + 1}"
        parts.append(f"<b>{_html.escape(label)}</b>\n<pre>{_html.escape(slide)}</pre>")
    parts.append(
        "\n💡 Используй фоны из Nana Banana (без текста), "
        "добавляй текст в Canva из Brand Kit."
    )
    return "\n\n".join(parts)


# ── Keyboards ────────────────────────────────────────────────────────────────

def _source_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([[
        InlineKeyboardButton("📈 На основе трендов", callback_data="ca:source:trends"),
        InlineKeyboardButton("✏️ Своя тема",          callback_data="ca:source:custom"),
    ]])


def _topics_keyboard(topics: list[str]) -> InlineKeyboardMarkup:
    row = [InlineKeyboardButton(str(i + 1), callback_data=f"ca:g:{i}") for i in range(len(topics))]
    buttons = [row[i:i + 5] for i in range(0, len(row), 5)]
    buttons.append([InlineKeyboardButton("🔄 Обновить темы", callback_data="ca:source:trends")])
    return InlineKeyboardMarkup(buttons)


def _action_buttons_no_images() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("🍌 Nana Banana (без текста)", callback_data="ca:prompt:notxt"),
            InlineKeyboardButton("🍌 Nana Banana (с текстом)",  callback_data="ca:prompt:text"),
        ],
        [
            InlineKeyboardButton("📄 PPTX (только тексты)", callback_data="ca:pptx:noimg"),
            InlineKeyboardButton("📋 Тексты для Canva",     callback_data="ca:prompt:canva"),
        ],
    ])


def _canva_buttons() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([[
        InlineKeyboardButton("📋 Тексты для Canva", callback_data="ca:prompt:canva"),
    ]])


def _pptx_from_my_images_button(count: int) -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([[
        InlineKeyboardButton(
            f"📄 Собрать PPTX ({count}/{len(_SLIDE_LABELS)} картинок)",
            callback_data="ca:pptx:userimages",
        )
    ]])


def _text_review_keyboard(n_slides: int) -> InlineKeyboardMarkup:
    """Keyboard for text-only review stage (before image generation)."""
    row = [InlineKeyboardButton(str(i + 1), callback_data=f"ca:edit:{i}") for i in range(n_slides)]
    buttons = [row[i:i + 5] for i in range(0, len(row), 5)]
    buttons.append([
        InlineKeyboardButton("🖼 Генерировать картинки", callback_data="ca:gen:images"),
        InlineKeyboardButton("🔄 Пересоздать",           callback_data="ca:regen:all"),
    ])
    buttons.append([
        InlineKeyboardButton("🍌 Промпты (с текстом)",  callback_data="ca:prompt:text"),
        InlineKeyboardButton("🍌 Промпты (фон)",        callback_data="ca:prompt:notxt"),
    ])
    return InlineKeyboardMarkup(buttons)


def _persist_carousel_draft(
    context: ContextTypes.DEFAULT_TYPE,
    topic: str,
    slides: list[str],
    img_prompts: list[str],
    angle: str = "",
    hook: str = "",
) -> str | None:
    from bot.services.drafts_store import save_draft as _save_draft, update_draft as _update_draft

    payload = {"slides": slides, "img_prompts": img_prompts}
    if hook:
        payload["hook"] = hook
    if angle:
        payload["angle"] = angle
    existing_draft_id = str(context.user_data.get("ca_draft_id", "")).strip()
    if existing_draft_id:
        updated = _update_draft(existing_draft_id, topic=topic, status="draft", payload=payload)
        if updated:
            context.user_data["ca_draft_id"] = updated.draft_id
            return updated.draft_id

    saved = _save_draft(
        kind="carousel",
        topic=topic,
        source="/carousel",
        payload=payload,
    )
    context.user_data["ca_draft_id"] = saved.draft_id
    return saved.draft_id


def _review_keyboard(n_slides: int, has_failed: bool = False) -> InlineKeyboardMarkup:
    """Keyboard for post-image review. Optionally shows retry-failed button."""
    row = [InlineKeyboardButton(str(i + 1), callback_data=f"ca:edit:{i}") for i in range(n_slides)]
    buttons = [row[i:i + 5] for i in range(0, len(row), 5)]
    action_row = [InlineKeyboardButton("📄 Скачать PPTX", callback_data="ca:pptx:final")]
    if has_failed:
        action_row.append(InlineKeyboardButton("🔄 Повторить ❌", callback_data="ca:regen:failed:note"))
    action_row.append(InlineKeyboardButton("🔄 Пересоздать всё", callback_data="ca:regen:all"))
    buttons.append(action_row)
    buttons.append([
        InlineKeyboardButton("🖼 Все с замечанием", callback_data="ca:regen:all:imgnote"),
    ])
    buttons.append([
        InlineKeyboardButton("🍌 Промпты (с текстом)", callback_data="ca:prompt:text"),
        InlineKeyboardButton("🍌 Промпты (фон)",       callback_data="ca:prompt:notxt"),
    ])
    return InlineKeyboardMarkup(buttons)


def _slide_actions_keyboard(idx: int) -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("🔄 Новый текст AI", callback_data=f"ca:edit:{idx}:ai"),
            InlineKeyboardButton("✏️ Свой текст",     callback_data=f"ca:edit:{idx}:manual"),
        ],
        [
            InlineKeyboardButton("🖼 Новая картинка",        callback_data=f"ca:edit:{idx}:img"),
            InlineKeyboardButton("🖼 С замечанием",          callback_data=f"ca:edit:{idx}:imgnote"),
        ],
        [
            InlineKeyboardButton("✅ Готово", callback_data="ca:review"),
        ],
    ])


# ── Slide-level regeneration ─────────────────────────────────────────────────

def _regen_slide_text_sync(topic: str, slides: list[str], idx: int) -> str:
    """Ask Claude to rewrite a single slide, aware of its role and neighbours."""
    from bot.services.claude_client import call_claude
    role = _SLIDE_VISUAL_ROLES[idx] if idx < len(_SLIDE_VISUAL_ROLES) else "supporting slide"
    label = _SLIDE_LABELS[idx] if idx < len(_SLIDE_LABELS) else f"Слайд {idx + 1}"
    others = "\n".join(
        f"Слайд {i + 1}: {s}" for i, s in enumerate(slides) if i != idx
    )
    prompt = (
        f"Ты — копирайтер карусели для Instagram. Ниша: регуляция нервной системы через сенсорные практики.\n\n"
        f"Тема карусели: {topic}\n"
        f"Роль {label}: {role}\n\n"
        f"Другие слайды (контекст):\n{others}\n\n"
        f"Напиши новый вариант текста для {label}.\n"
        f"Требования: максимум 5-6 строк, до 10 слов в строке, живой язык, без клише, без длинных тире.\n"
        f"Верни только текст слайда — ничего больше."
    )
    text = call_claude(
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300,
        context="carousel regen slide",
    )
    return _fix_dashes(text)


async def _show_slide_for_edit(
    target,
    idx: int,
    slides: list[str],
    images: list[bytes | None],
) -> None:
    """Send a single slide (image + text) with edit action buttons."""
    label = _SLIDE_LABELS[idx] if idx < len(_SLIDE_LABELS) else f"Слайд {idx + 1}"
    text = slides[idx]
    img = images[idx] if idx < len(images) else None
    caption = f"<b>{_html.escape(label)}</b>\n\n{_html.escape(text)}"

    if img:
        await target.reply_photo(
            photo=img,
            caption=caption,
            parse_mode="HTML",
            reply_markup=_slide_actions_keyboard(idx),
        )
    else:
        await target.reply_text(
            caption + "\n\n<i>картинка не сгенерирована</i>",
            parse_mode="HTML",
            reply_markup=_slide_actions_keyboard(idx),
        )


# ── Parallel image generation ────────────────────────────────────────────────

_FALLBACK_IMG_PROMPT = (
    "terracotta minimal lifestyle, incense smoke, soft natural light, "
    "--ar 4:5 --style atmospheric"
)


async def _run_image_generation(
    msg,
    context: ContextTypes.DEFAULT_TYPE,
    skip_existing: bool = False,
) -> None:
    """Generate slide images sequentially (Gemini has strict rate limits)."""
    slides      = context.user_data.get("ca_slides", [])
    img_prompts = context.user_data.get("ca_img_prompts", [])
    existing    = context.user_data.get("ca_gemini_images", [])
    n = len(slides)
    loop = asyncio.get_event_loop()

    # Start from existing images, extend to n slots
    images: list[bytes | None] = list(existing) + [None] * max(0, n - len(existing))

    # Which indices need generation
    indices = [
        i for i in range(n)
        if not (skip_existing and i < len(existing) and existing[i])
    ]

    if not indices:
        await msg.reply_text("✅ Все картинки уже готовы!")
        return

    processed: set[int] = set()
    # Pre-mark skipped slots as already processed
    for j in range(n):
        if j not in indices:
            processed.add(j)

    progress_msg = await msg.reply_text(
        f"🖼 Генерирую картинки: {'⏳' * n} 0/{len(indices)}"
    )
    done_count = 0

    async def gen_one(i: int) -> None:
        nonlocal done_count
        prompt = img_prompts[i] if i < len(img_prompts) else _FALLBACK_IMG_PROMPT
        img = await loop.run_in_executor(_img_executor, _gemini_slide, prompt, i)
        images[i] = img
        processed.add(i)
        done_count += 1
        icons = "".join(
            ("✅" if images[j] else "❌") if j in processed else "⏳"
            for j in range(n)
        )
        try:
            await progress_msg.edit_text(
                f"🖼 Картинки: {icons} {done_count}/{len(indices)}"
            )
        except Exception:
            pass
    await asyncio.gather(*[gen_one(i) for i in indices])

    try:
        await progress_msg.delete()
    except Exception:
        pass

    context.user_data["ca_gemini_images"] = images
    generated  = sum(1 for img in images if img)
    has_failed = generated < n

    # ── QA phase ─────────────────────────────────────────────────────────
    qa_results: dict[int, tuple[bool, str]] = {}
    generated_indices = [i for i, img in enumerate(images) if img]
    last_note = context.user_data.get("ca_last_note", "")

    if generated_indices and settings.anthropic_api_key:
        # ── QA round 1 ───────────────────────────────────────────────────
        qa_icons = ["➖"] * n
        for i in generated_indices:
            qa_icons[i] = "⏳"
        qa_progress = await msg.reply_text(
            f"🔍 Проверяю: {''.join(qa_icons)} 0/{len(generated_indices)}"
        )
        qa_done = 0

        async def qa_one(i: int) -> None:
            nonlocal qa_done
            prompt_i = img_prompts[i] if i < len(img_prompts) else ""
            passed, reason = await loop.run_in_executor(
                _executor, _qa_image_sync, images[i], prompt_i, last_note, i
            )
            qa_results[i] = (passed, reason)
            qa_done += 1
            qa_icons[i] = "✅" if passed else "⚠️"
            try:
                await qa_progress.edit_text(
                    f"🔍 Проверяю: {''.join(qa_icons)} {qa_done}/{len(generated_indices)}"
                )
            except Exception:
                pass

        await asyncio.gather(*[qa_one(i) for i in generated_indices])

        # ── Auto-rerender failed QA images ───────────────────────────────
        failed_qa = [i for i, (passed, _) in qa_results.items() if not passed]
        if failed_qa:
            for i in failed_qa:
                _, reason = qa_results[i]
                fix_note = f"fix: {reason}" if reason and reason.upper() != "OK" else "avoid impossible elements"
                old_prompt = img_prompts[i] if i < len(img_prompts) else _FALLBACK_IMG_PROMPT
                img_prompts[i] = _apply_note_to_prompt(old_prompt, fix_note)
                qa_icons[i] = "🔄"

            try:
                await qa_progress.edit_text(
                    f"🔄 Перерендер: {''.join(qa_icons)} — исправляю {len(failed_qa)} сл."
                )
            except Exception:
                pass

            regen_done = 0

            async def regen_one(i: int) -> None:
                nonlocal regen_done
                img = await loop.run_in_executor(_img_executor, _gemini_slide, img_prompts[i], i)
                images[i] = img
                regen_done += 1
                qa_icons[i] = "⏳" if img else "❌"
                try:
                    await qa_progress.edit_text(
                        f"🔄 Перерендер: {''.join(qa_icons)} {regen_done}/{len(failed_qa)}"
                    )
                except Exception:
                    pass

            await asyncio.gather(*[regen_one(i) for i in failed_qa])

            # ── QA round 2 ───────────────────────────────────────────────
            regenned = [i for i in failed_qa if images[i]]
            qa2_done = 0

            async def qa_two(i: int) -> None:
                nonlocal qa2_done
                prompt_i = img_prompts[i] if i < len(img_prompts) else ""
                passed, reason = await loop.run_in_executor(
                    _executor, _qa_image_sync, images[i], prompt_i, last_note, i
                )
                qa_results[i] = (passed, reason)
                qa2_done += 1
                qa_icons[i] = "✅" if passed else "⚠️"
                try:
                    await qa_progress.edit_text(
                        f"🔍 Повторная проверка: {''.join(qa_icons)} {qa2_done}/{len(regenned)}"
                    )
                except Exception:
                    pass

            if regenned:
                await asyncio.gather(*[qa_two(i) for i in regenned])

            context.user_data["ca_img_prompts"] = img_prompts
            context.user_data["ca_gemini_images"] = images

        try:
            await qa_progress.delete()
        except Exception:
            pass

    # Send all images in order
    for i, img in enumerate(images):
        if img:
            label = _SLIDE_LABELS[i] if i < len(_SLIDE_LABELS) else f"Слайд {i + 1}"
            caption = f"<b>{_html.escape(label)}</b>\n{_html.escape(slides[i])}"
            passed, reason = qa_results.get(i, (True, ""))
            if not passed and reason and reason.upper() != "OK":
                caption += f"\n\n⚠️ <i>{_html.escape(reason)}</i>"
            try:
                await msg.reply_photo(photo=img, caption=caption, parse_mode="HTML")
            except Exception:
                pass

    if not has_failed:
        pptx_bytes = await loop.run_in_executor(_executor, _build_pptx, slides, images)
        await msg.reply_document(
            document=io.BytesIO(pptx_bytes),
            filename="carousel.pptx",
            caption=(
                "📄 PPTX готов — фоны уже вставлены.\n"
                "Загрузи в Canva и настрой шрифт / цвета из Brand Kit."
            ),
        )
    else:
        await msg.reply_text(
            f"⚠️ Сгенерировано {generated}/{n} картинок. "
            "Нажми «🔄 Повторить ❌» чтобы попробовать ещё раз."
        )

    await msg.reply_text(
        "✏️ Нажми номер слайда чтобы изменить:",
        reply_markup=_review_keyboard(n, has_failed=has_failed),
    )


# ── Shared carousel generation helper ───────────────────────────────────────

async def _run_carousel(query_or_message, context: ContextTypes.DEFAULT_TYPE,
                        topic: str, status_msg) -> None:
    """Generate slide texts only. User then triggers image generation manually."""
    loop = asyncio.get_event_loop()

    await status_msg.edit_text(
        f"🎠 Тема: {topic}\n\n⏳ Генерирую черновик → прогоняю через редактора..."
    )

    slides, img_prompts, arc, angle, hook = await loop.run_in_executor(
        _executor, _generate_carousel_sync, topic
    )

    if not slides:
        target = query_or_message if hasattr(query_or_message, "reply_text") else query_or_message.message
        try:
            await status_msg.edit_text("❌ Не удалось сгенерировать карусель. Попробуй позже.")
        except Exception:
            await target.reply_text("❌ Не удалось сгенерировать карусель. Попробуй позже.")
        return

    context.user_data["ca_slides"]           = slides
    context.user_data["ca_img_prompts"]      = img_prompts
    context.user_data["ca_arc"]              = arc
    context.user_data["ca_topic"]            = topic
    context.user_data["ca_gemini_images"]    = []
    context.user_data["ca_awaiting_images"]  = False
    context.user_data["ca_user_image_ids"]   = []

    try:
        _persist_carousel_draft(context, topic, slides, img_prompts, angle=angle, hook=hook)
    except Exception:
        logger.exception("carousel: failed to save draft for topic: %s", topic)

    target = query_or_message if hasattr(query_or_message, "reply_text") else query_or_message.message

    keyboard = _text_review_keyboard(len(slides))

    lines = []
    for i, s in enumerate(slides):
        label = _SLIDE_LABELS[i] if i < len(_SLIDE_LABELS) else f"Слайд {i + 1}"
        lines.append(f"<b>{_html.escape(label)}</b>\n{_html.escape(s)}")
    slides_body = "\n\n".join(lines)

    header = "📝 <b>Тексты слайдов готовы:</b>\n\n"
    footer = "\n\nНажми номер слайда чтобы изменить, или генерируй картинки:"
    full_text = header + slides_body + footer

    # Telegram HTML limit is ~4096 rendered chars; split into two messages if needed
    if len(full_text) <= 4096:
        try:
            await status_msg.edit_text(full_text, parse_mode="HTML", reply_markup=keyboard)
        except Exception:
            try:
                await status_msg.delete()
            except Exception:
                pass
            await target.reply_text(full_text, parse_mode="HTML", reply_markup=keyboard)
    else:
        # Send first half of slides, then second half + keyboard
        mid = len(slides) // 2
        part1 = header + "\n\n".join(lines[:mid])
        part2 = "\n\n".join(lines[mid:]) + footer
        try:
            await status_msg.edit_text(part1, parse_mode="HTML")
        except Exception:
            try:
                await status_msg.delete()
            except Exception:
                pass
            await target.reply_text(part1, parse_mode="HTML")
        await target.reply_text(part2, parse_mode="HTML", reply_markup=keyboard)


# ── Command handler ──────────────────────────────────────────────────────────

async def cmd_carousel(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if not settings.anthropic_api_key:
        await update.message.reply_text("❌ Для /carousel нужен ANTHROPIC_API_KEY.")
        return

    context.user_data["ca_awaiting_images"]    = False
    context.user_data["ca_awaiting_topic"]     = False
    context.user_data["ca_awaiting_slide_edit"] = None
    context.user_data["ca_user_image_ids"]     = []

    await update.message.reply_text(
        "🎠 *Карусель для Instagram*\n\nВыбери, откуда взять тему:",
        parse_mode="Markdown",
        reply_markup=_source_keyboard(),
    )


# ── Callback handler ─────────────────────────────────────────────────────────

async def cb_carousel(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    from cache.store import cache
    from analytics.aggregator import collect_all

    query = update.callback_query
    await query.answer()
    data = query.data

    # ── Source: trends ────────────────────────────────────────────────────
    if data == "ca:source:trends":
        results = cache.get("results")
        if not results:
            await query.message.edit_text("⏳ Собираю тренды...")
            results = await collect_all()
            cache.set("results", results)

        await query.message.edit_text("🧠 Генерирую темы на основе трендов...")
        loop = asyncio.get_event_loop()
        topics = await loop.run_in_executor(
            _executor, _claude_topics_carousel, _format_trends(results)
        )

        if not topics:
            await query.message.edit_text("❌ Не удалось сгенерировать темы. Попробуй позже.")
            return

        context.user_data["ca_topics"] = topics
        items = "\n".join(f"{i+1}. {t}" for i, t in enumerate(topics))
        await query.message.edit_text(
            f"📈 Темы из трендов:\n\n{items}\n\nНажми номер — сгенерирую карусель:",
            reply_markup=_topics_keyboard(topics),
        )
        return

    # ── Source: custom topic ──────────────────────────────────────────────
    if data == "ca:source:custom":
        context.user_data["ca_awaiting_topic"] = True
        await query.message.edit_text(
            "✏️ Напиши тему для карусели одним сообщением.\n\n"
            "Например:\n"
            "— как запах помогает выйти из тревожной петли\n"
            "— 5 минут утром: сенсорный ритуал для старта\n"
            "— почему корпоративный wellbeing не работает без тела"
        )
        return

    # ── Pick topic from list ──────────────────────────────────────────────
    if data.startswith("ca:g:"):
        idx = int(data.split(":")[2])
        topics: list[str] = context.user_data.get("ca_topics", [])
        if not topics or idx >= len(topics):
            await query.message.reply_text("❌ Темы устарели — запроси /carousel снова.")
            return

        topic = topics[idx]
        status = await query.message.reply_text("⏳ Начинаю...")
        try:
            await _run_carousel(query, context, topic, status)
        except Exception:
            logger.exception("_run_carousel failed")
            await query.message.reply_text("❌ Ошибка при генерации. Попробуй ещё раз.")
        return

    # ── Prompt buttons ────────────────────────────────────────────────────
    if data in ("ca:prompt:text", "ca:prompt:notxt"):
        slides     = context.user_data.get("ca_slides", [])
        img_prompts = context.user_data.get("ca_img_prompts", [])
        topic      = context.user_data.get("ca_topic", "")
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        if not img_prompts:
            gen_msg = await query.message.reply_text("⏳ Генерирую промпты для картинок...")
            loop = asyncio.get_event_loop()
            img_prompts = await loop.run_in_executor(
                _executor, _generate_slide_image_prompts_sync, slides, topic
            )
            context.user_data["ca_img_prompts"] = img_prompts
            await gen_msg.delete()
        if data == "ca:prompt:text":
            await query.message.reply_text(
                _make_slide_prompts_with_text(img_prompts, slides),
                parse_mode="HTML",
            )
        else:
            await query.message.reply_text(
                _make_slide_prompts_no_text(img_prompts, slides),
                parse_mode="HTML",
            )
            context.user_data["ca_awaiting_images"] = True
            await query.message.reply_text(
                "📸 Сгенерировал в Nana Banana? Пришли картинки сюда — я соберу PPTX."
            )
        return

    if data == "ca:prompt:canva":
        slides = context.user_data.get("ca_slides", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        await query.message.reply_text(_format_for_canva(slides), parse_mode="HTML")
        return

    # ── Generate images (first time) ──────────────────────────────────────
    if data == "ca:gen:images":
        slides = context.user_data.get("ca_slides", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        await _run_image_generation(query.message, context, skip_existing=False)
        return

    # ── Regen ALL images with a note ─────────────────────────────────────
    if data == "ca:regen:all:imgnote":
        slides = context.user_data.get("ca_slides", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        context.user_data["ca_awaiting_img_note"] = {"idx": None, "skip_existing": False}
        await query.message.reply_text(
            "✏️ Напиши замечание — применю ко всем картинкам.\n"
            "<i>Например: более тёмные тона, без рук, добавить свечи, минималистичнее</i>",
            parse_mode="HTML",
        )
        return

    # ── Retry failed images — ask for note first ──────────────────────────
    if data == "ca:regen:failed:note":
        slides = context.user_data.get("ca_slides", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        context.user_data["ca_awaiting_img_note"] = {"idx": None, "skip_existing": True}
        await query.message.reply_text(
            "✏️ Напиши замечание к картинкам — что изменить, добавить или убрать.\n"
            "<i>Например: более тёмные тона, без рук, добавить свечи, минималистичнее</i>",
            parse_mode="HTML",
        )
        return

    # ── Retry failed images — no note (kept for backward compat) ─────────
    if data == "ca:regen:failed":
        slides = context.user_data.get("ca_slides", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        await _run_image_generation(query.message, context, skip_existing=True)
        return

    # ── Review screen ─────────────────────────────────────────────────────
    if data == "ca:review":
        slides = context.user_data.get("ca_slides", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        images = context.user_data.get("ca_gemini_images", [])
        if images:
            has_failed = any(img is None for img in images[:len(slides)])
            await query.message.reply_text(
                "✏️ Нажми номер слайда чтобы изменить:",
                reply_markup=_review_keyboard(len(slides), has_failed=has_failed),
            )
        else:
            await query.message.reply_text(
                "✏️ Нажми номер слайда чтобы изменить, или генерируй картинки:",
                reply_markup=_text_review_keyboard(len(slides)),
            )
        return

    # ── Regenerate whole carousel with same topic ──────────────────────────
    if data == "ca:regen:all":
        topic = context.user_data.get("ca_topic", "")
        if not topic:
            await query.message.reply_text("❌ Тема не найдена. Запроси /carousel заново.")
            return
        status = await query.message.reply_text("🔄 Пересоздаю карусель...")
        try:
            await _run_carousel(query, context, topic, status)
        except Exception:
            logger.exception("_run_carousel (regen) failed")
            await query.message.reply_text("❌ Ошибка при генерации. Попробуй ещё раз.")
        return

    # ── Slide editor ──────────────────────────────────────────────────────
    if data.startswith("ca:edit:"):
        parts = data.split(":")
        idx = int(parts[2])
        action = parts[3] if len(parts) > 3 else ""
        slides = context.user_data.get("ca_slides", [])
        images = context.user_data.get("ca_gemini_images", [])
        topic  = context.user_data.get("ca_topic", "")
        img_prompts = context.user_data.get("ca_img_prompts", [])

        if not slides or idx >= len(slides):
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return

        # ── Show slide ────────────────────────────────────────────────────
        if not action:
            await _show_slide_for_edit(query.message, idx, slides, images)
            return

        # ── AI: regenerate text (+ image) ─────────────────────────────────
        if action == "ai":
            status = await query.message.reply_text(
                f"🔄 Генерирую новый текст для слайда {idx + 1}..."
            )
            loop = asyncio.get_event_loop()
            new_text = await loop.run_in_executor(
                _executor, _regen_slide_text_sync, topic, slides, idx
            )
            slides[idx] = new_text
            context.user_data["ca_slides"] = slides

            if img_prompts and idx < len(img_prompts):
                await status.edit_text(f"🖼 Обновляю картинку для слайда {idx + 1}...")
                new_img = await loop.run_in_executor(
                    _img_executor, _gemini_slide, img_prompts[idx]
                )
                if new_img:
                    while len(images) <= idx:
                        images.append(None)
                    images[idx] = new_img
                    context.user_data["ca_gemini_images"] = images

            await status.delete()
            await _show_slide_for_edit(query.message, idx, slides, images)
            return

        # ── Manual: wait for user text ────────────────────────────────────
        if action == "manual":
            context.user_data["ca_awaiting_slide_edit"] = idx
            label = _SLIDE_LABELS[idx] if idx < len(_SLIDE_LABELS) else f"Слайд {idx + 1}"
            await query.message.reply_text(
                f"✏️ Введи новый текст для <b>{_html.escape(label)}</b>:\n"
                f"<i>Просто напиши следующим сообщением</i>",
                parse_mode="HTML",
            )
            return

        # ── Regenerate image only ─────────────────────────────────────────
        if action == "img":
            if not img_prompts or idx >= len(img_prompts):
                await query.message.reply_text("❌ Промт для картинки не найден.")
                return
            status = await query.message.reply_text(
                f"🖼 Генерирую новую картинку для слайда {idx + 1}..."
            )
            loop = asyncio.get_event_loop()
            new_img = await loop.run_in_executor(
                _img_executor, _gemini_slide, img_prompts[idx]
            )
            await status.delete()
            if new_img:
                while len(images) <= idx:
                    images.append(None)
                images[idx] = new_img
                context.user_data["ca_gemini_images"] = images
            else:
                await query.message.reply_text("⚠️ Gemini не сгенерировал картинку. Попробуй ещё раз.")
            await _show_slide_for_edit(query.message, idx, slides, images)
            return

        # ── Regenerate image with user note ──────────────────────────────
        if action == "imgnote":
            if not img_prompts or idx >= len(img_prompts):
                await query.message.reply_text("❌ Промт для картинки не найден.")
                return
            context.user_data["ca_awaiting_img_note"] = {"idx": idx, "skip_existing": False}
            label = _SLIDE_LABELS[idx] if idx < len(_SLIDE_LABELS) else f"Слайд {idx + 1}"
            await query.message.reply_text(
                f"✏️ Замечание для картинки <b>{_html.escape(label)}</b>:\n"
                "<i>Что изменить, добавить или убрать?</i>\n"
                "<i>Например: темнее, без рук, добавить свечи, более абстрактно</i>",
                parse_mode="HTML",
            )
            return

    # ── Final PPTX from current state ─────────────────────────────────────
    if data == "ca:pptx:final":
        slides = context.user_data.get("ca_slides", [])
        images = context.user_data.get("ca_gemini_images", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        status = await query.message.reply_text("⏳ Собираю PPTX...")
        loop = asyncio.get_event_loop()
        pptx_bytes = await loop.run_in_executor(_executor, _build_pptx, slides, images or None)
        await status.delete()
        await query.message.reply_document(
            document=io.BytesIO(pptx_bytes),
            filename="carousel_final.pptx",
            caption="📄 PPTX из текущей версии карусели.",
        )
        return

    # ── PPTX: text only ───────────────────────────────────────────────────
    if data == "ca:pptx:noimg":
        slides = context.user_data.get("ca_slides", [])
        if not slides:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return
        loop = asyncio.get_event_loop()
        pptx_bytes = await loop.run_in_executor(_executor, _build_pptx, slides, None)
        await query.message.reply_document(
            document=io.BytesIO(pptx_bytes),
            filename="carousel_texts.pptx",
            caption="📄 PPTX с текстами (без фонов). Замени фон в Canva на картинки из Nana Banana.",
        )
        return

    # ── PPTX: from user images ────────────────────────────────────────────
    if data == "ca:pptx:userimages":
        slides    = context.user_data.get("ca_slides", [])
        image_ids = context.user_data.get("ca_user_image_ids", [])
        if not slides or not image_ids:
            await query.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return

        status = await query.message.reply_text(
            f"⏳ Скачиваю {len(image_ids)} картинок и собираю PPTX..."
        )
        images: list[bytes | None] = []
        for file_id in image_ids[:len(_SLIDE_LABELS)]:
            try:
                tg_file = await context.bot.get_file(file_id)
                buf = await tg_file.download_as_bytearray()
                images.append(bytes(buf))
            except Exception as exc:
                logger.warning("Failed to download user image: %s", exc)
                images.append(None)

        loop = asyncio.get_event_loop()
        pptx_bytes = await loop.run_in_executor(_executor, _build_pptx, slides, images)
        await status.delete()
        await query.message.reply_document(
            document=io.BytesIO(pptx_bytes),
            filename="carousel_with_images.pptx",
            caption="📄 PPTX с твоими картинками готов. Загрузи в Canva и настрой шрифты / цвета.",
        )
        context.user_data["ca_awaiting_images"] = False
        return


# ── Message handlers ──────────────────────────────────────────────────────────

async def msg_carousel_topic(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle text input — either a new topic or manual slide edit."""
    text = (update.message.text or "").strip()

    # ── Image note received ────────────────────────────────────────────────
    img_note_state = context.user_data.get("ca_awaiting_img_note")
    if img_note_state is not None:
        if not text:
            return
        context.user_data["ca_awaiting_img_note"] = None
        context.user_data["ca_last_note"] = text
        slides     = context.user_data.get("ca_slides", [])
        img_prompts = list(context.user_data.get("ca_img_prompts", []))
        images     = context.user_data.get("ca_gemini_images", [])
        slide_idx_note  = img_note_state.get("idx")        # None = all failed
        skip_existing   = img_note_state.get("skip_existing", False)

        if not slides:
            await update.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            return

        if slide_idx_note is None:
            # Apply note to all prompts (for retry-failed flow)
            img_prompts = [_apply_note_to_prompt(p, text) for p in img_prompts]
        else:
            # Apply note to single slide prompt
            if slide_idx_note < len(img_prompts):
                img_prompts[slide_idx_note] = _apply_note_to_prompt(
                    img_prompts[slide_idx_note], text
                )

        context.user_data["ca_img_prompts"] = img_prompts

        if slide_idx_note is None:
            await _run_image_generation(update.message, context, skip_existing=skip_existing)
        else:
            status = await update.message.reply_text(
                f"🖼 Генерирую картинку для слайда {slide_idx_note + 1} с замечанием..."
            )
            loop = asyncio.get_event_loop()
            new_img = await loop.run_in_executor(
                _img_executor, _gemini_slide, img_prompts[slide_idx_note]
            )
            await status.delete()
            if new_img:
                while len(images) <= slide_idx_note:
                    images.append(None)
                images[slide_idx_note] = new_img
                context.user_data["ca_gemini_images"] = images
            else:
                await update.message.reply_text("⚠️ Gemini не сгенерировал картинку. Попробуй ещё раз.")
            await _show_slide_for_edit(update.message, slide_idx_note, slides, images)
        return

    # ── Manual slide edit ──────────────────────────────────────────────────
    slide_idx = context.user_data.get("ca_awaiting_slide_edit")
    if slide_idx is not None:
        if not text:
            return
        slides = context.user_data.get("ca_slides", [])
        images = context.user_data.get("ca_gemini_images", [])
        img_prompts = context.user_data.get("ca_img_prompts", [])

        if not slides or slide_idx >= len(slides):
            await update.message.reply_text("❌ Данные устарели. Сгенерируй карусель заново.")
            context.user_data["ca_awaiting_slide_edit"] = None
            return

        slides[slide_idx] = text
        context.user_data["ca_slides"] = slides
        context.user_data["ca_awaiting_slide_edit"] = None

        # Regenerate image for the edited slide
        if img_prompts and slide_idx < len(img_prompts):
            status = await update.message.reply_text("✅ Текст обновлён! 🖼 Генерирую картинку...")
            loop = asyncio.get_event_loop()
            new_img = await loop.run_in_executor(
                _executor, _gemini_slide, img_prompts[slide_idx]
            )
            await status.delete()
            if new_img:
                while len(images) <= slide_idx:
                    images.append(None)
                images[slide_idx] = new_img
                context.user_data["ca_gemini_images"] = images

        await _show_slide_for_edit(update.message, slide_idx, slides, images)
        return

    # ── New carousel topic ─────────────────────────────────────────────────
    if not context.user_data.get("ca_awaiting_topic"):
        return

    if len(text) < 5:
        await update.message.reply_text("❌ Тема слишком короткая. Опиши подробнее.")
        return

    context.user_data["ca_awaiting_topic"] = False
    status = await update.message.reply_text("⏳ Начинаю...")
    try:
        await _run_carousel(update.message, context, text, status)
    except Exception:
        logger.exception("_run_carousel (topic msg) failed")
        await update.message.reply_text("❌ Ошибка при генерации. Попробуй ещё раз.")


async def msg_carousel_photo(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Collect user-sent photos for PPTX assembly."""
    if not context.user_data.get("ca_awaiting_images"):
        return

    photo = update.message.photo[-1]
    ids: list[str] = context.user_data.setdefault("ca_user_image_ids", [])
    max_slides = len(_SLIDE_LABELS)

    if len(ids) >= max_slides:
        await update.message.reply_text(
            f"✅ Уже {max_slides} картинок — нажми кнопку для сборки PPTX."
        )
        return

    ids.append(photo.file_id)
    count = len(ids)

    await update.message.reply_text(
        f"✅ Получено {count}/{max_slides}."
        + ("" if count < max_slides else " Все получены!") +
        " Пришли ещё или собирай PPTX:" if count < max_slides else " Собирай PPTX:",
        reply_markup=_pptx_from_my_images_button(count),
    )


def build_carousel_handler():
    return [
        CommandHandler("carousel", cmd_carousel),
        CallbackQueryHandler(cb_carousel, pattern="^ca:"),
        MessageHandler(filters.PHOTO, msg_carousel_photo),
        MessageHandler(filters.TEXT & ~filters.COMMAND, msg_carousel_topic),
    ]
